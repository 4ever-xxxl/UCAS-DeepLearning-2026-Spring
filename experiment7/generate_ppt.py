"""
Generate experiment7 presentation PPT based on the report outline.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
OUTPUT_DIR = Path(__file__).resolve().parent / "outputs"
PPT_PATH = OUTPUT_DIR / "experiment7_report.pptx"
METRICS_PATH = OUTPUT_DIR / "metrics.json"

# ---------------------------------------------------------------------------
# Theme colors
# ---------------------------------------------------------------------------
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)
SECONDARY = RGBColor(0x2D, 0x3A, 0x4A)
ACCENT = RGBColor(0x00, 0x96, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFA)
DARK_TEXT = RGBColor(0x1E, 0x1E, 0x1E)
CODE_BG = RGBColor(0x28, 0x2C, 0x34)
CODE_FG = RGBColor(0xAB, 0xB2, 0xBF)
TABLE_HEADER_BG = RGBColor(0x1A, 0x56, 0xDB)
TABLE_ROW_ALT = RGBColor(0xE8, 0xEE, 0xF8)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

SLIDE_W = Cm(33.867)  # 16:9
SLIDE_H = Cm(19.05)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Cm(left), Cm(top), Cm(width), Cm(height),
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_bullet_list(tf, items, font_size=14, color=DARK_TEXT, bullet_char="●"):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(6)


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = add_rect(slide, left, top, width, height, fill_color=CODE_BG)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.4)
    tf.margin_right = Cm(0.4)
    tf.margin_top = Cm(0.3)
    tf.margin_bottom = Cm(0.3)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = CODE_FG
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.LEFT


def add_image(slide, path, left, top, width, height=None):
    if height is None:
        slide.shapes.add_picture(str(path), Cm(left), Cm(top), Cm(width))
    else:
        slide.shapes.add_picture(str(path), Cm(left), Cm(top), Cm(width), Cm(height))


def add_page_number(slide, num, total):
    add_textbox(slide, 30, 17.8, 3.5, 0.8, f"{num} / {total}",
                font_size=10, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 33.867, 1.0, fill_color=PRIMARY)
    add_textbox(slide, 1.0, 0.1, 30, 0.9, title, font_size=22, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 1.0, 1.15, 30, 0.8, subtitle, font_size=12, color=RGBColor(0x88, 0x99, 0xAA))


def add_table(slide, left, top, headers, rows, col_widths=None, font_size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Cm(left), Cm(top),
                                         Cm(sum(col_widths or [6] * n_cols)),
                                         Cm(0.8 * n_rows))
    tbl = table_shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER

    return table_shape


def _load_metrics():
    return json.loads(METRICS_PATH.read_text())


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build():
    m = _load_metrics()

    best_val_ppl = m["best_validation_perplexity"]
    best_val_loss = m["best_validation_loss"]
    test_loss = m["test_loss"]
    test_ppl = m["test_perplexity"]
    test_acc = m["test_next_word_accuracy"] * 100
    meets_req = m["meets_requirement"]
    vocab_size = m["vocab_size"]
    param_count = m["parameter_count"]
    history = m["history"]

    epoch_1 = history[0]
    epoch_10 = history[9]
    epoch_20 = history[19]
    epoch_55 = history[-1]

    TOTAL_PAGES = 11

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # =========================================================================
    # Page 1: Cover
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)

    add_rect(slide, 0, 0, 33.867, 0.3, fill_color=PRIMARY)
    add_rect(slide, 0, 3.5, 33.867, 9.0, fill_color=LIGHT_BG)

    add_textbox(slide, 2.0, 4.5, 30, 2.0,
                "实验七：神经网络语言模型",
                font_size=36, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 2.0, 7.0, 30, 1.5,
                "Penn Treebank 上的 LSTM 词级语言模型",
                font_size=20, color=SECONDARY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 2.0, 13.5, 30, 0.8,
                "深度学习 · 2026 Spring", font_size=14, color=RGBColor(0x88, 0x99, 0xAA),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 2.0, 14.5, 30, 0.8,
                "李相廷 · 2026-05-11",
                font_size=12, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 1, TOTAL_PAGES)

    # =========================================================================
    # Page 2: TOC
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)
    add_section_header(slide, "目录", "Contents")

    toc_items = [
        "1. 实验任务概述",
        "2. 数据集介绍",
        "3. 模型结构设计",
        "4. 损失函数与评价指标",
        "5. 优化器与训练策略",
        "6. 工程实现与创新点",
        "7. 训练配置",
        "8. 实验结果",
        "9. 结果分析与讨论",
        "10. 总结与改进方向",
    ]

    for i, item in enumerate(toc_items):
        y = 2.5 + i * 1.35
        add_rect(slide, 4.0, y, 0.4, 0.4, fill_color=PRIMARY)
        add_textbox(slide, 5.0, y - 0.1, 25, 1.0, item, font_size=18, color=SECONDARY)

    add_page_number(slide, 2, TOTAL_PAGES)

    # =========================================================================
    # Page 3: Task Overview
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)
    add_section_header(slide, "实验任务概述", "任务定义 | 应用场景 | 实验目标")

    add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "任务定义", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 1.2, 3.0, 15.0, 5.0, "", font_size=13)
    add_bullet_list(tf, [
        "词级语言模型（Word-level Language Model）",
        "给定前文词序列，预测下一个词的条件概率分布",
        "核心公式：P(w₁, ..., wₙ) = ∏ P(wₜ | w₁, ..., wₜ₋₁)",
        "属于 NLP 基础任务，是众多下游应用的基石",
        "使用截断 BPTT 保持跨 batch 的隐藏状态连续性",
    ], font_size=13)

    add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "应用场景", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 17.7, 3.0, 15.0, 5.0, "", font_size=13)
    add_bullet_list(tf, [
        "机器翻译中的候选译文重排序",
        "语音识别的语言模型打分",
        "文本生成与自动补全",
        "预训练语言模型的基础组件（ELMo、GPT 等）",
        "拼写纠错与输入法",
    ], font_size=13)

    add_rect(slide, 1.0, 8.5, 31.5, 0.6, fill_color=ACCENT)
    add_textbox(slide, 1.2, 8.5, 31, 0.6, "实验目标", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 1.2, 9.4, 31, 8.0, "", font_size=13)
    add_bullet_list(tf, [
        "在 Penn Treebank (PTB) 语料库上训练一个 LSTM 词级语言模型",
        "核心指标：测试集困惑度（Perplexity, PPL）低于 80",
        "辅助指标：next-word top-1 准确率",
        "技术路线：Embedding + 2层 LSTM + Dropout + Softmax，截断 BPTT 训练",
    ], font_size=13)

    add_page_number(slide, 3, TOTAL_PAGES)

    # =========================================================================
    # Page 4: Dataset
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)
    add_section_header(slide, "数据集介绍", "Penn Treebank (PTB) simple-examples")

    add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "数据集概况", font_size=16, bold=True, color=WHITE)

    add_table(slide, 1.2, 3.2,
              ["属性", "详情"],
              [
                  ["数据集", "Penn Treebank (PTB) simple-examples"],
                  ["词表大小", f"{vocab_size:,}"],
                  ["训练集", "约 93 万 token"],
                  ["验证集", "约 7.4 万 token"],
                  ["测试集", "约 8.2 万 token"],
                  ["特殊 token", "<eos> 句末 · <unk> 低频/未知词"],
              ],
              col_widths=[7, 8])

    add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "数据预处理", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 17.7, 3.2, 15, 6, "", font_size=12)
    add_bullet_list(tf, [
        "换行符 → <eos>，保留句子边界信息",
        "仅用训练集统计词频、构建 10,000 词表",
        "valid / test 集低频词统一映射为 <unk>",
        "文本按 batch_size × batch_len reshape 为连续 token 流",
        "训练时取 num_steps=35 的滑动窗口",
        "目标序列 = 输入序列右移一位",
    ], font_size=12)

    add_rect(slide, 1.0, 8.6, 31.5, 0.6, fill_color=ACCENT)
    add_textbox(slide, 1.2, 8.6, 31, 0.6, "批数据生成示意", font_size=14, bold=True, color=WHITE)

    add_code_block(slide, 1.2, 9.5, 31, 2.5,
                   "# 连续文本 → reshape → 滑动窗口\n"
                   "# raw: [w1 w2 w3 ... wN]\n"
                   "# reshape: batch_size × batch_len 矩阵\n"
                   "# 每个 step 取 num_steps 列作为输入\n"
                   "# inputs[b, t] = data[b, step*num_steps + t]\n"
                   "# targets[b, t] = data[b, step*num_steps + t + 1]",
                   font_size=11)

    add_page_number(slide, 4, TOTAL_PAGES)

    # =========================================================================
    # Page 5: Model Architecture
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)
    add_section_header(slide, "模型结构设计", "两层 LSTM 语言模型 · 66M 参数")

    add_rect(slide, 1.0, 2.2, 31.5, 2.6, fill_color=LIGHT_BG)
    add_textbox(slide, 1.5, 2.3, 30.5, 0.6, "数据流", font_size=14, bold=True, color=PRIMARY)
    add_textbox(slide, 2.0, 3.0, 30, 1.5,
                "word_ids [B, T]\n"
                "    → Embedding(10000, 1500) → Dropout(0.65)\n"
                "    → LSTM(1500, 1500, num_layers=2, dropout=0.65)\n"
                "    → Dropout(0.65) → Linear(1500, 10000)\n"
                "    → logits [B×T, 10000] → CrossEntropyLoss",
                font_size=13, color=SECONDARY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.5, 5.2, 31, 0.6, "large preset 结构参数", font_size=14, bold=True, color=PRIMARY)
    add_table(slide, 1.5, 5.9,
              ["模块", "配置", "说明"],
              [
                  ["Embedding", "10,000 × 1,500", "词向量维度 1500"],
                  ["LSTM", "2 层, hidden=1500", "层间 dropout=0.65"],
                  ["Dropout", "0.65", "嵌入层后 + LSTM 输出后"],
                  ["Classifier", "Linear 1500→10000", "输出词表概率分布"],
                  ["Weight init", "Uniform(-0.05, 0.05)", "所有参数均匀初始化"],
                  ["总参数量", f"{param_count:,}", "66.03 M"],
              ],
              col_widths=[7, 8, 13])

    add_textbox(slide, 1.5, 11.5, 31, 0.6, "设计理由", font_size=14, bold=True, color=PRIMARY)
    tf = add_textbox(slide, 1.5, 12.2, 31, 5, "", font_size=12)
    add_bullet_list(tf, [
        "LSTM 门控机制有效缓解长序列梯度消失，适合顺序文本建模",
        "大 hidden size=1500 提供充足的模型容量",
        "强 dropout=0.65 是针对 PTB 小语料（~1M token）防止过拟合的关键设计",
        "不采用 weight tying：保持 embedding 和输出投影独立，给模型更多自由度",
    ], font_size=12)

    add_page_number(slide, 5, TOTAL_PAGES)

    # =========================================================================
    # Page 6: Loss & Metrics
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)
    add_section_header(slide, "损失函数与评价指标", "CrossEntropyLoss | Perplexity | Next-Word Accuracy")

    add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "损失函数", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 1.2, 3.0, 15, 6, "", font_size=12)
    add_bullet_list(tf, [
        "CrossEntropyLoss（交叉熵损失）",
        "输入：模型输出 logits [B×T, V]",
        "目标：真实下一个词 ID [B×T]",
        "等价于最大化真实词的条件概率",
        "loss = -1/N Σ log P(wₜ | context)",
    ], font_size=12)

    add_code_block(slide, 1.2, 7.5, 15, 1.5,
                   "criterion = nn.CrossEntropyLoss()\n"
                   "logits, hidden = model(inputs, hidden)\n"
                   "loss = criterion(logits, targets.reshape(-1))",
                   font_size=10)

    add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=ACCENT)
    add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "评价指标", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 17.7, 3.0, 15, 6, "", font_size=12)
    add_bullet_list(tf, [
        "主要指标：困惑度 Perplexity (PPL)",
        "  PPL = exp(cross_entropy_loss)",
        "  PPL 越低 → 模型对测试文本建模越好",
        "  实验要求：测试集 PPL < 80",
        "辅助指标：next-word top-1 accuracy",
        "  预测概率最高的词是否与真实词一致",
        "  PPL 比 accuracy 更全面（考虑概率分布）",
    ], font_size=12)

    add_rect(slide, 1.0, 9.8, 31.5, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 1.2, 9.8, 31, 0.6, "PPL 与 Accuracy 的关系", font_size=14, bold=True, color=WHITE)

    tf = add_textbox(slide, 1.2, 10.6, 31, 5, "", font_size=12)
    add_bullet_list(tf, [
        "词表有 10,000 类，很多上下文存在多个合理的下一个词",
        "Accuracy 只看 top-1 预测是否正确，过于严格",
        "PPL 衡量模型对整个概率分布的质量，更能反映语言建模能力",
        "例如：上下文 \"the cat sat on the ___\" 中 \"mat\" 和 \"floor\" 都合理，accuracy 只认一个",
    ], font_size=12)

    add_page_number(slide, 6, TOTAL_PAGES)

    # =========================================================================
    # Page 7: Optimizer & Training Strategy
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)
    add_section_header(slide, "优化器与训练策略", "SGD · 学习率衰减 · 梯度裁剪 · AMP")

    # Optimizer
    add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "优化器配置", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 1.2, 3.0, 15, 6, "", font_size=12)
    add_bullet_list(tf, [
        "优化器：SGD（随机梯度下降）",
        "初始学习率：1.0",
        "学习率衰减：第 15 轮起，每轮 × 1/1.15",
        "  lrₜ = lrₜ₋₁ / 1.15  (t ≥ 15)",
        "梯度裁剪：global norm ≤ 10.0",
        "  防止 BPTT 中梯度爆炸",
        "Epochs：55",
    ], font_size=12)

    add_code_block(slide, 1.2, 8.5, 15, 2.0,
                   "optimizer = torch.optim.SGD(\n"
                   "    model.parameters(), lr=1.0\n"
                   ")\n"
                   "clip_grad_norm_(model.parameters(), 10.0)\n"
                   "if epoch >= 15:\n"
                   "    lr /= 1.15",
                   font_size=10)

    # Strategy
    add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=ACCENT)
    add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "训练策略", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 17.7, 3.0, 15, 6, "", font_size=12)
    add_bullet_list(tf, [
        "截断 BPTT：num_steps=35",
        "  每 35 步截断梯度，隐藏状态跨 batch 传递",
        "Batch size：20",
        "  保持 PTB 原始顺序，不 shuffle",
        "AMP 混合精度（bfloat16）",
        "TF32 加速（Tensor Core）",
        "torch.compile(mode=\"reduce-overhead\")",
        "硬件：NVIDIA RTX 5070 (12GB)",
    ], font_size=12)

    add_rect(slide, 1.0, 11.2, 31.5, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 1.2, 11.2, 31, 0.6, "学习率衰减曲线", font_size=14, bold=True, color=WHITE)

    add_table(slide, 1.2, 12.0,
              ["Epoch", "1", "10", "14", "15", "20", "30", "40", "55"],
              [["LR", "1.000", "1.000", "1.000", "0.870", "0.432", "0.107", "0.026", "0.003"]],
              col_widths=[3] + [2.8] * 8)

    add_page_number(slide, 7, TOTAL_PAGES)

    # =========================================================================
    # Page 8: Engineering & Innovations
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)
    add_section_header(slide, "工程实现与创新点", "PyTorch 复现 · Preset 机制 · 鲁棒数据下载")

    add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "工程实现亮点", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 1.2, 3.0, 15, 7, "", font_size=12)
    add_bullet_list(tf, [
        "完整 PyTorch 复现 TensorFlow PTB 经典流程",
        "PTBBatchDataset 对齐原始 ptb_producer：",
        "  连续文本流、固定 batch、目标右移一位",
        "stateful BPTT：隐藏状态跨 batch 保持",
        "  每个 epoch 开始时 detach 隐藏状态",
        "可复用 preset 机制：",
        "  test → 快速验证代码链路",
        "  small / medium → 调试超参",
        "  large → 正式达标训练",
    ], font_size=12)

    add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=ACCENT)
    add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "创新点", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 17.7, 3.0, 15, 7, "", font_size=12)
    add_bullet_list(tf, [
        "鲁棒数据下载：",
        "  首选 Mikolov simple-examples.tgz",
        "  失败时回退 PTB 原始三文件镜像",
        "  下载后 MD5 校验完整性",
        "现代训练优化：",
        "  AMP + TF32 + torch.compile",
        "  在保持经典架构的前提下大幅提升效率",
        "丰富输出：",
        "  best checkpoint + metrics.json",
        "  训练曲线 + next-word 样例 + 续写样例",
    ], font_size=12)

    add_rect(slide, 1.0, 10.5, 31.5, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 1.2, 10.5, 31, 0.6, "Preset 配置一览", font_size=14, bold=True, color=WHITE)

    add_table(slide, 1.2, 11.3,
              ["Preset", "Embedding", "Hidden", "Layers", "Dropout", "Epochs", "Batch", "用途"],
              [
                  ["test", "200", "200", "2", "0.5", "3", "20", "快速验证链路"],
                  ["small", "200", "200", "2", "0.5", "20", "20", "轻量调试"],
                  ["medium", "650", "650", "2", "0.5", "40", "20", "中等容量"],
                  ["large", "1500", "1500", "2", "0.65", "55", "20", "达标训练"],
              ],
              col_widths=[3, 3.5, 3, 2.5, 2.5, 2.5, 2.5, 6])

    add_page_number(slide, 8, TOTAL_PAGES)

    # =========================================================================
    # Page 9: Training Configuration
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)
    add_section_header(slide, "训练配置", "Training Configuration")

    add_textbox(slide, 1.2, 2.2, 31, 0.6, "正式训练命令", font_size=14, bold=True, color=PRIMARY)
    add_code_block(slide, 1.2, 2.9, 31, 1.2,
                   "uv run python experiment7/train_ptb_lm.py --preset large --device cuda "
                   "--output-dir experiment7/outputs",
                   font_size=11)

    add_textbox(slide, 1.2, 4.5, 31, 0.6, "超参数汇总", font_size=14, bold=True, color=PRIMARY)
    add_table(slide, 1.2, 5.2,
              ["参数", "值", "参数", "值"],
              [
                  ["Preset", "large", "Epochs", "55"],
                  ["Batch Size", "20", "BPTT Steps", "35"],
                  ["Embedding Dim", "1500", "Hidden Dim", "1500"],
                  ["LSTM Layers", "2", "Dropout", "0.65"],
                  ["Optimizer", "SGD", "Initial LR", "1.0"],
                  ["LR Decay", "÷1.15/epoch (≥15)", "Grad Clip", "10.0"],
                  ["AMP", "bfloat16", "TF32", "Enabled"],
                  ["torch.compile", "reduce-overhead", "Hardware", "RTX 5070 12GB"],
              ],
              col_widths=[7, 8, 7, 8])

    add_textbox(slide, 1.2, 12.2, 31, 0.6, "训练流程", font_size=14, bold=True, color=PRIMARY)
    tf = add_textbox(slide, 1.2, 12.9, 31, 5, "", font_size=12)
    add_bullet_list(tf, [
        "Step 1: 下载 PTB 数据 → MD5 校验 → 构建词表 → 生成 batch 数据",
        "Step 2: 构建 PTBLanguageModel (66M params) → 部署到 CUDA",
        "Step 3: 55 epoch 训练/验证循环，每 epoch 记录 loss、ppl、accuracy、lr",
        "Step 4: 自动保存验证 PPL 最优的 checkpoint",
        "Step 5: 加载最优模型 → 测试集评估 → 输出 metrics.json",
        "Step 6: 运行 generate_results.py → 训练曲线 + next-word 样例 + 续写样例",
    ], font_size=12)

    add_page_number(slide, 9, TOTAL_PAGES)

    # =========================================================================
    # Page 10: Experimental Results
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)
    add_section_header(slide, "实验结果", "关键指标 | 训练曲线")

    # Key metrics highlight
    add_rect(slide, 1.0, 2.2, 31.5, 1.5, fill_color=LIGHT_BG)
    add_textbox(slide, 1.5, 2.3, 6.5, 1.3,
                f"最佳验证 PPL\n{best_val_ppl:.2f}",
                font_size=22, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 9.0, 2.3, 6.5, 1.3,
                f"测试 PPL\n{test_ppl:.2f}",
                font_size=22, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 16.5, 2.3, 6.5, 1.3,
                f"测试 Loss\n{test_loss:.4f}",
                font_size=22, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 24.0, 2.3, 6.5, 1.3,
                f"Next-Word Acc\n{test_acc:.2f}%",
                font_size=18, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

    # Goal check
    add_rect(slide, 1.0, 4.0, 31.5, 0.8, fill_color=GREEN if meets_req else RGBColor(0xE7, 0x4C, 0x3C))
    status_text = "✓ 达到实验要求：测试集 PPL < 80" if meets_req else "✗ 未达到实验要求"
    add_textbox(slide, 1.2, 4.0, 31, 0.8, status_text,
                font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Training progress table
    add_textbox(slide, 1.5, 5.2, 31, 0.6, "训练过程摘要", font_size=14, bold=True, color=PRIMARY)
    add_table(slide, 1.5, 5.9,
              ["Epoch", "1", "5", "10", "15", "20", "30", "40", "55"],
              [
                  ["LR", "1.000", "1.000", "1.000", "0.870", "0.432", "0.107", "0.026", "0.003"],
                  ["Train PPL", f"{epoch_1['train_perplexity']:.1f}", "138.1",
                   "102.0", "83.3", "61.9", "44.9", "40.5", "39.2"],
                  ["Val PPL", f"{epoch_1['val_perplexity']:.1f}", "129.8",
                   "103.5", "93.5", "85.3", "82.8", "82.7", "82.6"],
              ],
              col_widths=[3] + [2.8] * 8, font_size=10)

    # Training curves image
    curves_path = OUTPUT_DIR / "training_curves.png"
    if curves_path.exists():
        add_textbox(slide, 1.5, 8.8, 31, 0.6, "训练曲线（Loss / Perplexity / Accuracy / Learning Rate）", font_size=14, bold=True, color=PRIMARY)
        add_image(slide, curves_path, 1.5, 9.4, 31, 7.5)

    add_page_number(slide, 10, TOTAL_PAGES)

    # =========================================================================
    # Page 11: Analysis & Summary
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, WHITE)
    add_section_header(slide, "结果分析与总结", "关键发现 | 总结 | 改进方向")

    # Analysis
    add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "关键发现", font_size=14, bold=True, color=WHITE)
    tf = add_textbox(slide, 1.2, 3.0, 15, 5, "", font_size=11)
    add_bullet_list(tf, [
        "训练初期 PPL 下降极快：epoch 1 的 548 →",
        "  epoch 10 的 102，模型快速学到高频词",
        "  和 <eos> 句末模式",
        "Dropout=0.65 对大容量 LSTM 至关重要：",
        "  小语料下防止过拟合的核心手段",
        "梯度裁剪保证 BPTT 中 SGD 更新稳定",
        "学习率衰减是后期突破瓶颈的关键：",
        "  epoch 15 后衰减使 val PPL 从 93 降至 82",
    ], font_size=11)

    add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=ACCENT)
    add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "PPL vs Accuracy", font_size=14, bold=True, color=WHITE)
    tf = add_textbox(slide, 17.7, 3.0, 15, 5, "", font_size=11)
    add_bullet_list(tf, [
        "测试 PPL=78.94 达标，但 accuracy 仅 27.86%",
        "原因：10,000 类词表的 top-1 预测本身极难",
        "PPL 评估整个概率分布的质量",
        "  比单一标签正确性更全面",
        "模型在后期验证 PPL 趋于平缓",
        "  接近该架构在 PTB 上的泛化瓶颈",
        "训练 PPL (39.2) 与验证 PPL (82.6) 差距大",
        "  存在过拟合，可增强正则化",
    ], font_size=11)

    # Summary
    add_rect(slide, 1.0, 8.3, 31.5, 0.6, fill_color=PRIMARY)
    add_textbox(slide, 1.2, 8.3, 31, 0.6, "完成情况", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 1.5, 9.1, 31, 3, "", font_size=12)
    add_bullet_list(tf, [
        "完成 PTB 词级 LSTM 语言模型的完整实现：数据下载与校验、词表构建、BPTT 训练、评估、可视化",
        f"large preset 测试集 PPL={test_ppl:.2f}，达到实验要求的 PPL<80",
        f"模型参数量 {param_count:,}，训练 55 epoch 后在验证集上取得最佳 PPL={best_val_ppl:.2f}",
    ], font_size=12)

    # Improvements
    add_rect(slide, 1.0, 12.0, 31.5, 0.6, fill_color=ACCENT)
    add_textbox(slide, 1.2, 12.0, 31, 0.6, "改进方向", font_size=16, bold=True, color=WHITE)

    tf = add_textbox(slide, 1.5, 12.8, 15, 5, "", font_size=11)
    add_bullet_list(tf, [
        "Weight Tying：共享 embedding 和输出投影权重",
        "  减少参数、可能提升泛化",
        "AWD-LSTM 正则化：variational dropout +",
        "  weight dropout + AR/TAR 正则化",
        "  预期将 PPL 降至 60 以下",
    ], font_size=11)

    tf = add_textbox(slide, 17.5, 12.8, 15, 5, "", font_size=11)
    add_bullet_list(tf, [
        "Transformer 对比：替换为 Decoder-only",
        "  Transformer，对比两种架构差异",
        "解码策略：top-k / top-p 采样",
        "  替代 argmax 贪心，提升文本续写多样性",
        "更大语料：WikiText-2 / WikiText-103",
        "  测试模型在更大语料上的表现",
    ], font_size=11)

    add_page_number(slide, 11, TOTAL_PAGES)

    # -----------------------------------------------------------------------
    # Save
    # -----------------------------------------------------------------------
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPT_PATH))
    print(f"PPT saved to: {PPT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
