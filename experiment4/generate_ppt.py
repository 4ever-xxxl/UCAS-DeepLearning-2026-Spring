"""
Generate experiment4 presentation PPT based on the report outline.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
OUTPUT_DIR = Path(__file__).resolve().parent / "outputs"
PPT_PATH = OUTPUT_DIR / "experiment4_report.pptx"
METRICS_PATH = OUTPUT_DIR / "metrics.json"

# ---------------------------------------------------------------------------
# Theme colors
# ---------------------------------------------------------------------------
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)
SECONDARY = RGBColor(0x2D, 0x3A, 0x4A)
ACCENT = RGBColor(0x00, 0x96, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFA)
DARK_TEXT = RGBColor(0x1E, 0x1E, 0x1E)
CODE_BG = RGBColor(0x28, 0x2C, 0x34)
CODE_FG = RGBColor(0xAB, 0xB2, 0xBF)
TABLE_HEADER_BG = RGBColor(0x1A, 0x56, 0xDB)
TABLE_ROW_ALT = RGBColor(0xE8, 0xEE, 0xF8)

SLIDE_W = Cm(33.867)  # 16:9
SLIDE_H = Cm(19.05)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Cm(left), Cm(top), Cm(width), Cm(height),
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_bullet_list(tf, items, font_size=14, color=DARK_TEXT, bullet_char="●"):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(6)


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = add_rect(slide, left, top, width, height, fill_color=CODE_BG)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.4)
    tf.margin_right = Cm(0.4)
    tf.margin_top = Cm(0.3)
    tf.margin_bottom = Cm(0.3)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = CODE_FG
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.LEFT


def add_image(slide, path, left, top, width, height=None):
    if height is None:
        slide.shapes.add_picture(str(path), Cm(left), Cm(top), Cm(width))
    else:
        slide.shapes.add_picture(str(path), Cm(left), Cm(top), Cm(width), Cm(height))


def add_page_number(slide, num, total):
    add_textbox(slide, 30, 17.8, 3.5, 0.8, f"{num} / {total}",
                font_size=10, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 33.867, 1.0, fill_color=PRIMARY)
    add_textbox(slide, 1.0, 0.1, 30, 0.9, title, font_size=22, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 1.0, 1.15, 30, 0.8, subtitle, font_size=12, color=RGBColor(0x88, 0x99, 0xAA))


def add_table(slide, left, top, headers, rows, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Cm(left), Cm(top),
                                         Cm(sum(col_widths or [6] * n_cols)),
                                         Cm(0.8 * n_rows))
    tbl = table_shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER

    return table_shape


# ---------------------------------------------------------------------------
# Load metrics
# ---------------------------------------------------------------------------
metrics = json.loads(METRICS_PATH.read_text())
best_val_bleu = metrics["best_val_bleu"]
test_loss = metrics["test_loss"]
test_ppl = metrics["test_ppl"]
test_bleu = metrics["test_bleu"]
test_bleu_beam = metrics["test_bleu_beam"]

TOTAL_PAGES = 10

# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

# ===========================================================================
# Page 1: 封面
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)

add_rect(slide, 0, 0, 33.867, 0.3, fill_color=PRIMARY)

add_rect(slide, 0, 3.5, 33.867, 9.0, fill_color=LIGHT_BG)
add_textbox(slide, 2.0, 4.5, 30, 2.0,
            "实验四：基于 Transformer 的中英神经机器翻译",
            font_size=36, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 2.0, 7.0, 30, 1.5,
            "从零实现 Encoder–Decoder + BPE 子词分词",
            font_size=20, color=SECONDARY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 2.0, 13.5, 30, 0.8,
            "深度学习 · 2026 Spring", font_size=14, color=RGBColor(0x88, 0x99, 0xAA),
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, 2.0, 14.5, 30, 0.8,
            "李相廷 · 2026-05-10",
            font_size=12, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

add_page_number(slide, 1, TOTAL_PAGES)

# ===========================================================================
# Page 2: 目录
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "目录", "Contents")

toc_items = [
    "1. 任务概述",
    "2. 数据集介绍",
    "3. Transformer 模型结构",
    "4. 损失函数与优化器",
    "5. 训练过程与配置",
    "6. 实验结果",
    "7. 结果分析与讨论",
    "8. 总结与改进方向",
]

for i, item in enumerate(toc_items):
    y = 2.5 + i * 1.5
    add_rect(slide, 4.0, y, 0.4, 0.4, fill_color=PRIMARY)
    add_textbox(slide, 5.0, y - 0.1, 25, 1.0, item, font_size=18, color=SECONDARY)

add_page_number(slide, 2, TOTAL_PAGES)

# ===========================================================================
# Page 3: 任务概述
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "任务概述", "任务定义 | 应用场景 | 实验目标")

add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "任务定义", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 3.0, 15.0, 5.0, "", font_size=13)
add_bullet_list(tf, [
    "中文 → 英文神经机器翻译（Neural Machine Translation）",
    "输入：中文句子（已 ICTCLAS 分词，空格分隔）",
    "输出：英文句子（tokenized + lower-case）",
    "属于典型的 seq2seq 序列生成任务",
    "核心挑战：变长输入输出、长距离语义依赖、OOV 问题",
], font_size=13)

add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "应用场景", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 17.7, 3.0, 15.0, 5.0, "", font_size=13)
add_bullet_list(tf, [
    "跨语言信息检索与机器翻译",
    "跨境电商产品描述翻译",
    "国际新闻自动翻译",
    "多语言智能客服系统",
], font_size=13)

add_rect(slide, 1.0, 8.5, 31.5, 0.6, fill_color=ACCENT)
add_textbox(slide, 1.2, 8.5, 31, 0.6, "实验目标", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 9.4, 31, 8.0, "", font_size=13)
add_bullet_list(tf, [
    "从零实现 Transformer Encoder–Decoder 架构（Self-Attention、Cross-Attention、Positional Encoding 等）",
    "使用 SentencePiece BPE 子词分词解决 OOV 问题，完成数据预处理与词汇表构建",
    "采用 Teacher Forcing 训练 + Warmup-Rsqrt 学习率调度 + Beam Search 解码",
    "在 NiuTrans 测试集上达到 BLEU-4 ≥ 14 的翻译质量",
], font_size=13)

add_page_number(slide, 3, TOTAL_PAGES)

# ===========================================================================
# Page 4: 数据集介绍
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "数据集介绍", "NiuTrans 中英平行语料库详情与 BPE 分词")

add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "数据集详情", font_size=16, bold=True, color=WHITE)

add_table(slide, 1.2, 3.2,
          ["属性", "详情"],
          [
              ["来源", "NiuTrans（东北大学 NLP Lab）"],
              ["训练集", "99,000 句中英平行句对"],
              ["验证集", "1,000 句对"],
              ["测试集", "1,000 句中文，各含英文参考"],
              ["中文", "ICTCLAS 分词，空格分隔 ~26 tok/句"],
              ["英文", "已 tokenize + lower-case ~30 tok/句"],
          ],
          col_widths=[7, 8])

add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "BPE 子词分词", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 17.7, 3.2, 15, 6, "", font_size=12)
add_bullet_list(tf, [
    "双 SentencePiece BPE 模型：zh & en 各 16,000 词",
    "char_coverage: zh=0.9995, en=1.0",
    "OOV 率近乎 0%（词级别曾有 ~38-51% UNK）",
    "特殊 token：<pad>=0, <unk>=1, <s>=2, </s>=3",
    "编码后目标序列前后加 <s> 和 </s>",
    "动态填充对齐 batch 内序列长度",
], font_size=12)

add_rect(slide, 1.0, 8.6, 31.5, 0.6, fill_color=ACCENT)
add_textbox(slide, 1.2, 8.6, 31, 0.6, "数据预处理流水线修复", font_size=14, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 9.4, 31, 8, "", font_size=12)
add_bullet_list(tf, [
    "原始数据流水线存在两个 bug：dev.zh / dev.en 行错位 + test.en 内容误置为中文",
    "通过 _parse_paired_blocks 正确解析 NiuTrans 三行块格式（zh ⏎ blank ⏎ en）",
    "训练集 / 验证集从 100k TM 数据中随机切分（seed=0，dev_size=1000）",
    "训练时每个 batch 动态 padding 到该 batch 最大长度（bucketing by length）",
], font_size=12)

add_page_number(slide, 4, TOTAL_PAGES)

# ===========================================================================
# Page 5: Transformer 模型结构
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "Transformer 模型结构", "Encoder–Decoder 架构 · Pre-LayerNorm · 19.25M 参数")

add_rect(slide, 1.0, 2.2, 31.5, 2.6, fill_color=LIGHT_BG)
add_textbox(slide, 1.5, 2.3, 30.5, 0.6, "整体流程", font_size=14, bold=True, color=PRIMARY)
add_textbox(slide, 2.0, 3.0, 30, 1.5,
            "src_ids → Embed(√d) + PE → [EncoderLayer × 6] → enc_norm\n"
            "                                                          ↘\n"
            "tgt_ids → Embed(√d) + PE → [DecoderLayer × 6] → dec_norm → output_proj → softmax\n"
            "         (tied with output_proj-weight sharing)",
            font_size=13, color=SECONDARY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 5.2, 31, 0.6, "模型超参与结构参数", font_size=14, bold=True, color=PRIMARY)
add_table(slide, 1.5, 5.9,
          ["模块 / 超参", "配置", "说明"],
          [
              ["d_model", "256", "模型隐藏维度"],
              ["n_heads", "8", "多头注意力头数"],
              ["Encoder layers", "6", "编码器层数"],
              ["Decoder layers", "6", "解码器层数"],
              ["d_ff", "1024", "FFN 中间层维度"],
              ["Dropout", "0.1", "全模型统一"],
              ["Weight Tying", "✓", "目标词嵌入 = 输出投影权重"],
              ["总参数量", "19,252,224", "19.25 M"],
          ],
          col_widths=[7, 8, 13])

add_textbox(slide, 1.5, 11.5, 31, 0.6, "设计关键点", font_size=14, bold=True, color=PRIMARY)
tf = add_textbox(slide, 1.5, 12.2, 31, 5, "", font_size=12)
add_bullet_list(tf, [
    "Pre-LayerNorm 残差：x + Dropout(SubLayer(LayerNorm(x)))，训练更稳定",
    "Multi-Head Self-Attention / Cross-Attention / Position-wise FFN + 正余弦位置编码",
    "使用 F.scaled_dot_product_attention（Flash Attention），mask 语义：True=attend",
    "调试关键：与 nn.MultiheadAttention 的 mask 语义相反（后者 True=mask），修正后 BLEU 从 0 跃升",
], font_size=12)

add_page_number(slide, 5, TOTAL_PAGES)

# ===========================================================================
# Page 6: 损失函数与优化器
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "损失函数与优化器", "Label Smoothing | AdamW | Warmup-Rsqrt | AMP")

# Loss
add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "损失函数", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 3.0, 15, 6, "", font_size=12)
add_bullet_list(tf, [
    "CrossEntropyLoss + Label Smoothing（ε=0.1）",
    "Label Smoothing：把 one-hot 目标 y 软化为 (1-ε)·y + ε/V",
    "有效抑制模型对训练分布的过度拟合",
    "提升 BLEU 指标（Vaswani et al. 2017）",
    "忽略 PAD 位置（ignore_index=PAD_IDX）",
], font_size=12)

add_code_block(slide, 1.2, 8.5, 15, 1.5,
                "criterion = nn.CrossEntropyLoss(\n"
                "    ignore_index=PAD_IDX,\n"
                "    label_smoothing=0.1\n"
                ")",
                font_size=10)

# Optimizer
add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "优化器设计", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 17.7, 3.0, 15, 5, "", font_size=12)
add_bullet_list(tf, [
    "AdamW 优化器",
    "β₁=0.9, β₂=0.98, ε=1e-9",
    "weight_decay = 1e-5（L2 正则化）",
    "Warmup-Rsqrt 学习率调度",
    "warmup_steps = 2000（~2.6 epoch）",
    "peak lr = 7e-4",
    "梯度裁剪 clip_grad_norm_=1.0",
], font_size=12)

# LR formula
add_code_block(slide, 17.5, 8.0, 15, 2.5,
                "# warmup-rsqrt schedule\n"
                "if step < warmup_steps:\n"
                "    lr = peak_lr * step/warmup\n"
                "else:\n"
                "    lr = peak_lr * sqrt(warmup/step)",
                font_size=10)

# AMP
add_rect(slide, 1.0, 10.8, 31.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 10.8, 31, 0.6, "AMP 混合精度训练", font_size=14, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 11.6, 31, 5, "", font_size=12)
add_bullet_list(tf, [
    "CUDA 环境下启用 bfloat16 AMP（torch.amp.autocast + GradScaler）",
    "RTX 5070 上 bfloat16 无精度损失，训练效率显著提升",
    "通过验证 BLEU 早停，仅保存验证集最优 checkpoint（best_model.pt）",
    "Decoder 训练使用 Teacher Forcing，评估时采用 Greedy / Beam Search 解码",
], font_size=12)

add_page_number(slide, 6, TOTAL_PAGES)

# ===========================================================================
# Page 7: 训练过程与配置
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "训练过程与配置", "Training Configuration & Pipeline")

add_textbox(slide, 1.2, 2.2, 31, 0.6, "训练超参数", font_size=14, bold=True, color=PRIMARY)
add_table(slide, 1.2, 2.9,
          ["参数", "值", "参数", "值"],
          [
              ["Batch Size", "128", "Max Epochs", "30"],
              ["Peak LR", "7e-4", "Warmup Steps", "2,000"],
              ["Optimizer", "AdamW", "Loss", "CrossEntropy + LS(0.1)"],
              ["Grad Clip", "1.0", "Dropout", "0.1"],
              ["AMP", "bfloat16", "Seed", "0"],
              ["Max Seq Len", "128 BPE tokens", "Beam Size", "5 (lp=0.6)"],
              ["Hardware", "NVIDIA RTX 5070", "Framework", "PyTorch 2.11"],
              ["Train Time", "~30 min", "Best Epoch", "28 / 30"],
          ],
          col_widths=[7, 8, 7, 8])

add_textbox(slide, 1.2, 10.2, 31, 0.6, "训练流程", font_size=14, bold=True, color=PRIMARY)
tf = add_textbox(slide, 1.2, 10.9, 31, 7, "", font_size=12)
add_bullet_list(tf, [
    "Step 1：修复原始数据流水线 bug，训练双 BPE 模型，编码源 / 目标语言",
    "Step 2：构建 TransformerNMT 模型（19.25M 参数），部署到 CUDA 设备",
    "Step 3：30 轮 Teacher Forcing 训练 / 验证，每轮记录 loss、ppl、bleu、lr",
    "Step 4：自动保存验证 BLEU 最高的模型 checkpoint（epoch 28, BLEU 26.69）",
    "Step 5：加载最优模型，在独立测试集上评估 Greedy 和 Beam Search 解码",
    "Step 6：输出 metrics.json 和 test_samples.json（含完整训练历史与翻译样例）",
], font_size=12)

# Key epoch milestones
add_textbox(slide, 1.2, 15.8, 31, 0.6, "关键 BLEU 增长节点", font_size=13, bold=True, color=ACCENT)
add_table(slide, 1.2, 16.5,
          ["Epoch", "1", "4", "8", "20", "25", "28 (Best)", "30"],
          [["Val BLEU", "1.01", "16.72", "23.08", "25.31", "26.63", "26.69", "26.48"]],
          col_widths=[3] + [2.8] * 7)

add_page_number(slide, 7, TOTAL_PAGES)

# ===========================================================================
# Page 8: 实验结果
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "实验结果", "关键指标 | 训练曲线")

# Key metrics highlight
add_rect(slide, 1.0, 2.2, 31.5, 1.5, fill_color=LIGHT_BG)
add_textbox(slide, 1.5, 2.3, 6.5, 1.3,
            f"最佳验证 BLEU\n{best_val_bleu:.2f}",
            font_size=22, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 9.0, 2.3, 6.5, 1.3,
            f"测试 BLEU (Greedy)\n{test_bleu:.2f}",
            font_size=22, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 16.5, 2.3, 6.5, 1.3,
            f"测试 BLEU (Beam=5)\n{test_bleu_beam:.2f}",
            font_size=22, bold=True, color=RGBColor(0xE6, 0x7E, 0x22),
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, 24.0, 2.3, 6.5, 1.3,
            f"Test Loss / PPL\n{test_loss:.2f} / {test_ppl:.2f}",
            font_size=18, bold=True, color=RGBColor(0x27, 0xAE, 0x60),
            alignment=PP_ALIGN.CENTER)

# Training curves
add_textbox(slide, 1.5, 4.2, 31, 0.6, "训练曲线（Loss / Perplexity / BLEU / Learning Rate）", font_size=14, bold=True, color=PRIMARY)
add_image(slide, OUTPUT_DIR / "training_curves.png", 1.5, 4.8, 31, 11.5)

add_page_number(slide, 8, TOTAL_PAGES)

# ===========================================================================
# Page 9: 结果分析与讨论
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "结果分析与讨论", "关键发现 | 翻译样例 | 误差分析")

# Finding 1
add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "Mask 语义修复是关键转折", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 1.2, 3.0, 15, 3, "", font_size=12)
add_bullet_list(tf, [
    "F.scaled_dot_product_attention 使用 True=attend 约定",
    "与 nn.MultiheadAttention 的 True=mask 语义相反",
    "修正后 epoch 1 BLEU=1.0 → epoch 4 BLEU=16.7",
    "目标值在 4 个 epoch 内即被超越",
], font_size=12)

# Finding 2
add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "BPE 子词分词的必要性", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 17.7, 3.0, 15, 3, "", font_size=12)
add_bullet_list(tf, [
    "词级别词表产生 ~38-51% UNK 率，BLEU ≈ 0",
    "BPE 把 vocab 限定 16k，OOV 率近乎 0%",
    "中英文各训练独立 BPE 模型",
    "子词粒度平衡了词表大小与语义完整性",
], font_size=12)

# Finding 3
add_rect(slide, 1.0, 6.0, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 6.0, 14.5, 0.6, "Beam Search 与收敛趋势", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 1.2, 6.8, 15, 3, "", font_size=12)
add_bullet_list(tf, [
    "Beam=5 + length penalty=0.6 比 greedy 提升 ~0.93 BLEU",
    "符合标准 NMT 经验，beam search 缓解曝光偏差",
    "Epoch 20 已达 25.3，再 10 epoch 仅多 ~1 BLEU",
    "19M 参数 + 100k 句对 下呈现收益递减",
], font_size=12)

# Finding 4
add_rect(slide, 17.5, 6.0, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 17.7, 6.0, 14.5, 0.6, "典型翻译样例", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 17.7, 6.8, 15, 3, "", font_size=11)
add_bullet_list(tf, [
    "✅ \"second , comprehensive management .\"（完全一致）",
    "✅ 突尼斯赞扬中国政策的长句（语义一致，语序略调）",
    "⚠️ 出现 \"branch branch branches\" token 重复",
    "❌ 专有名词 / 长难句有漏译和错译",
], font_size=11)

# Error analysis
add_rect(slide, 1.0, 10.5, 31.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 10.5, 31, 0.6, "误差分析与改进启示", font_size=14, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 11.3, 15, 5, "", font_size=12)
add_bullet_list(tf, [
    "重复 token：解码时可加入 no-repeat-ngram 启发式",
    "专有名词错译/漏译：训练数据中长尾命名实体不足",
    "长难句信息丢失：模型容量（d=256）不足以捕获全部语义",
], font_size=12)

tf = add_textbox(slide, 17.5, 11.3, 15, 5, "", font_size=12)
add_bullet_list(tf, [
    "增大模型（d_model=512, ~65M）可望突破 30 BLEU",
    "回译（Back-Translation）扩充训练数据",
    "用 mBART / NLLB 等预训练模型微调",
], font_size=12)

add_page_number(slide, 9, TOTAL_PAGES)

# ===========================================================================
# Page 10: 总结与改进方向
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "总结与改进方向", "Summary & Future Work")

# Summary
add_rect(slide, 1.0, 2.2, 31.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 31, 0.6, "完成情况", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.5, 3.0, 31, 4, "", font_size=13)
add_bullet_list(tf, [
    "从零实现了完整 Transformer Encoder–Decoder 架构（Self-Attn / Cross-Attn / PE / Pre-LN 残差）",
    "实现了双 SentencePiece BPE 子词分词流水线，解决了词级别高 OOV 导致 BLEU=0 的问题",
    "修复了原始数据流水线的两个 bug（dev 行错位 + test.en 内容错误），保证评估可靠性",
    f"Test BLEU 27.79（beam=5）/ 26.86（greedy），远超实验目标 BLEU ≥ 14，接近 2 倍",
], font_size=13)

# Improvements
add_rect(slide, 1.0, 7.5, 31.5, 0.6, fill_color=ACCENT)
add_textbox(slide, 1.2, 7.5, 31, 0.6, "改进方向", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.5, 8.3, 15, 8, "", font_size=12)
add_bullet_list(tf, [
    "模型规模升级",
    "  - Transformer-base: d_model=512, ~65M params",
    "  - 期望 BLEU 提升至 ~30",
    "数据增广",
    "  - 回译（Back-Translation）扩充训练数据",
    "  - 域内单语数据利用",
], font_size=12)

tf = add_textbox(slide, 17.5, 8.3, 15, 8, "", font_size=12)
add_bullet_list(tf, [
    "解码增强",
    "  - no-repeat-ngram 防止 token 重复",
    "  - coverage penalty / ensemble 解码",
    "预训练模型迁移",
    "  - mBART / NLLB 微调",
    "  - 更大规模预训练知识迁移",
    "评估增强",
    "  - 多参考 BLEU 评估（当前仅 1 参考）",
], font_size=12)

add_page_number(slide, 10, TOTAL_PAGES)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(PPT_PATH))
print(f"PPT saved to: {PPT_PATH}")
print(f"Total slides: {len(prs.slides)}")
