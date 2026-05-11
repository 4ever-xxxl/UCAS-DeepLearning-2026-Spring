"""
Generate experiment2 presentation PPT based on the report outline.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
OUTPUT_DIR = Path(__file__).resolve().parent / "outputs"
PPT_PATH = OUTPUT_DIR / "experiment2_report.pptx"
METRICS_PATH = OUTPUT_DIR / "metrics.json"

# ---------------------------------------------------------------------------
# Theme colors
# ---------------------------------------------------------------------------
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)  # Deep blue
SECONDARY = RGBColor(0x2D, 0x3A, 0x4A)  # Dark gray
ACCENT = RGBColor(0x00, 0x96, 0x88)  # Teal
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFA)
DARK_TEXT = RGBColor(0x1E, 0x1E, 0x1E)
CODE_BG = RGBColor(0x28, 0x2C, 0x34)
CODE_FG = RGBColor(0xAB, 0xB2, 0xBF)
TABLE_HEADER_BG = RGBColor(0x1A, 0x56, 0xDB)
TABLE_ROW_ALT = RGBColor(0xE8, 0xEE, 0xF8)

SLIDE_W = Cm(33.867)  # 16:9
SLIDE_H = Cm(19.05)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    bold=False,
    color=DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    font_name="Microsoft YaHei",
):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Cm(left),
        Cm(top),
        Cm(width),
        Cm(height),
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_bullet_list(tf, items, font_size=14, color=DARK_TEXT, bullet_char="●"):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(6)


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    """Add a code block with dark background."""
    shape = add_rect(slide, left, top, width, height, fill_color=CODE_BG)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.4)
    tf.margin_right = Cm(0.4)
    tf.margin_top = Cm(0.3)
    tf.margin_bottom = Cm(0.3)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = CODE_FG
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.LEFT


def add_image(slide, path, left, top, width, height=None):
    if height is None:
        slide.shapes.add_picture(str(path), Cm(left), Cm(top), Cm(width))
    else:
        slide.shapes.add_picture(str(path), Cm(left), Cm(top), Cm(width), Cm(height))


def add_page_number(slide, num, total):
    add_textbox(
        slide,
        30,
        17.8,
        3.5,
        0.8,
        f"{num} / {total}",
        font_size=10,
        color=RGBColor(0x99, 0xAA, 0xBB),
        alignment=PP_ALIGN.RIGHT,
    )


def add_section_header(slide, title, subtitle=""):
    """Add a colored header bar at the top of a content slide."""
    add_rect(slide, 0, 0, 33.867, 1.0, fill_color=PRIMARY)
    add_textbox(slide, 1.0, 0.1, 30, 0.9, title, font_size=22, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 1.0, 1.15, 30, 0.8, subtitle, font_size=12, color=RGBColor(0x88, 0x99, 0xAA))


def add_table(slide, left, top, headers, rows, col_widths=None):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Cm(left), Cm(top), Cm(sum(col_widths or [6] * n_cols)), Cm(0.8 * n_rows)
    )
    tbl = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER

    return table_shape


# ---------------------------------------------------------------------------
# Load metrics
# ---------------------------------------------------------------------------
metrics = json.loads(METRICS_PATH.read_text())
best_val_acc = metrics["best_validation_accuracy"]
test_acc = metrics["test_accuracy"]
test_loss = metrics["test_loss"]

TOTAL_PAGES = 10

# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

# ===========================================================================
# Page 1: Title Slide
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)

# Top accent bar
add_rect(slide, 0, 0, 33.867, 0.3, fill_color=PRIMARY)

# Title area
add_rect(slide, 0, 3.5, 33.867, 9.0, fill_color=LIGHT_BG)
add_textbox(
    slide,
    2.0,
    4.5,
    30,
    2.0,
    "实验二：基于ViT的CIFAR-10图像分类",
    font_size=36,
    bold=True,
    color=PRIMARY,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    2.0,
    7.0,
    30,
    1.5,
    "Vision Transformer 从零实现与训练",
    font_size=20,
    color=SECONDARY,
    alignment=PP_ALIGN.CENTER,
)

# Bottom info
add_textbox(
    slide,
    2.0,
    13.5,
    30,
    0.8,
    "深度学习 · 2026 Spring",
    font_size=14,
    color=RGBColor(0x88, 0x99, 0xAA),
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    2.0,
    14.5,
    30,
    0.8,
    "PyTorch 2.11 · NVIDIA RTX 5070 · CIFAR-10",
    font_size=12,
    color=RGBColor(0x99, 0xAA, 0xBB),
    alignment=PP_ALIGN.CENTER,
)

add_page_number(slide, 1, TOTAL_PAGES)

# ===========================================================================
# Page 2: 概述 Overview
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "概述", "任务目标 | 数据集 | 解决方案")

# Left column - Task
add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "任务目标", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 3.0, 15.0, 4.5, "", font_size=13)
add_bullet_list(
    tf,
    [
        "使用 PyTorch 从零实现 Vision Transformer (ViT) 图像分类模型",
        "深入理解 Attention 机制、Transformer Encoder 和 MLP 分类头",
        "掌握图像分类任务的完整流程：数据加载→模型构建→训练→验证→测试",
        "在 CIFAR-10 测试集上达到 80% 以上分类准确率",
    ],
    font_size=13,
)

# Right column - Dataset
add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "数据集概况", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 17.7, 3.0, 15.0, 4.5, "", font_size=13)
add_bullet_list(
    tf,
    [
        "CIFAR-10：60,000 张 32×32 彩色图像",
        "训练集 50,000 张 / 测试集 10,000 张",
        "10 个类别：airplane, automobile, bird, cat, deer, dog, frog, horse, ship, truck",
        "实验划分：训练集 45,000 + 验证集 5,000 + 测试集 10,000",
    ],
    font_size=13,
)

# Solution overview
add_rect(slide, 1.0, 8.0, 31.5, 0.6, fill_color=ACCENT)
add_textbox(slide, 1.2, 8.0, 31, 0.6, "解决方案概览", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 8.9, 31, 8.0, "", font_size=13)
add_bullet_list(
    tf,
    [
        "模型架构：ConvPatchEmbed → CLS Token + Position Embedding → 6× Transformer Encoder → LayerNorm → Linear(10)",
        "轻量化设计：使用卷积 tokenizer 直接处理 32×32 输入，避免放大到 224×224，减少计算量并增强局部特征提取",
        "训练策略：AdamW 优化器 + Warmup + Cosine Annealing + 混合精度训练（AMP）",
        "正则化增强：MixUp、RandAugment、RandomErasing、Label Smoothing、Dropout、Stochastic Depth",
        "结果：最佳验证准确率 84.60%，测试准确率 83.36%，满足 80% 指标要求",
    ],
    font_size=13,
)

add_page_number(slide, 2, TOTAL_PAGES)

# ===========================================================================
# Page 3: 数据集介绍
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "数据集介绍", "CIFAR-10 数据集详情与预处理")

# Left - dataset details
add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "数据集详情", font_size=16, bold=True, color=WHITE)

add_table(
    slide,
    1.2,
    3.2,
    ["属性", "值"],
    [
        ["图像尺寸", "3 × 32 × 32"],
        ["训练集规模", "50,000 张"],
        ["测试集规模", "10,000 张"],
        ["类别数", "10"],
        ["训练:验证划分", "9:1 (45,000:5,000)"],
    ],
    col_widths=[7, 8],
)

add_textbox(slide, 1.2, 8.2, 15, 0.6, "10 个类别", font_size=14, bold=True, color=SECONDARY)
tf = add_textbox(slide, 1.2, 8.8, 15, 4, "", font_size=12)
add_bullet_list(
    tf,
    [
        "airplane · automobile · bird · cat · deer",
        "dog · frog · horse · ship · truck",
    ],
)

# Right - preprocessing
add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "数据预处理与增强", font_size=16, bold=True, color=WHITE)

add_textbox(slide, 17.7, 3.2, 14.5, 0.6, "标准化", font_size=13, bold=True, color=ACCENT)
tf = add_textbox(slide, 17.7, 3.8, 14.5, 1.5, "", font_size=11)
add_bullet_list(
    tf,
    [
        "mean = (0.4914, 0.4822, 0.4465)",
        "std = (0.2023, 0.1994, 0.2010)",
    ],
    font_size=11,
    bullet_char="·",
)

add_textbox(slide, 17.7, 5.5, 14.5, 0.6, "训练数据增强", font_size=13, bold=True, color=ACCENT)
tf = add_textbox(slide, 17.7, 6.1, 14.5, 5, "", font_size=11)
add_bullet_list(
    tf,
    [
        "RandomCrop(32, padding=4) — 随机裁剪",
        "RandomHorizontalFlip — 随机水平翻转",
        "RandAugment(num_ops=2, magnitude=9) — 自动增强",
        "RandomErasing(p=0.15) — 随机擦除",
    ],
    font_size=11,
)

add_textbox(slide, 17.7, 10.5, 14.5, 0.6, "正则化策略", font_size=13, bold=True, color=ACCENT)
tf = add_textbox(slide, 17.7, 11.1, 14.5, 5, "", font_size=11)
add_bullet_list(
    tf,
    [
        "MixUp (α=0.2) — 混合样本增强",
        "Label Smoothing (ε=0.05) — 标签平滑",
        "Dropout (p=0.1) — 随机失活",
        "Stochastic Depth (p=0.1) — 随机深度",
    ],
    font_size=11,
)

add_page_number(slide, 3, TOTAL_PAGES)

# ===========================================================================
# Page 4: 网络结构设计
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "网络结构设计", "ViT 整体架构与核心模块")

# Architecture diagram (text-based)
add_rect(slide, 1.0, 2.2, 31.5, 2.2, fill_color=LIGHT_BG)
add_textbox(slide, 1.5, 2.3, 30.5, 0.6, "整体流程", font_size=14, bold=True, color=PRIMARY)
add_textbox(
    slide,
    2.0,
    3.0,
    30,
    1.2,
    "Input(3×32×32)  →  ConvPatchEmbed  →  [CLS] Token + Pos Embed  →  "
    "6× Transformer Encoder Blocks  →  LayerNorm  →  Linear(10)",
    font_size=14,
    color=SECONDARY,
    alignment=PP_ALIGN.CENTER,
)

# Model hyperparams table
add_textbox(slide, 1.5, 4.8, 31, 0.6, "模型超参数", font_size=14, bold=True, color=PRIMARY)
add_table(
    slide,
    1.5,
    5.5,
    ["参数", "值", "参数", "值"],
    [
        ["Embed Dim", "256", "Depth", "6"],
        ["Num Heads", "8", "MLP Ratio", "3.0"],
        ["Patch方式", "Conv Stem (4×)", "Drop Path Rate", "0.1"],
        ["Dropout", "0.1", "QKV Bias", "True"],
    ],
    col_widths=[7, 8, 7, 8],
)

# Conv vs Standard PatchEmbed comparison
add_textbox(slide, 1.5, 9.8, 31, 0.6, "ConvPatchEmbed：卷积 Tokenizer 设计", font_size=14, bold=True, color=PRIMARY)

add_code_block(
    slide,
    1.5,
    10.5,
    15,
    7.5,
    "ConvPatchEmbed (32×32 input):\n"
    "  Conv2d(3→64, k3, s1, p1) + BN + GELU\n"
    "  Conv2d(64→128, k3, s2, p1) + BN + GELU\n"
    "  Conv2d(128→256, k3, s2, p1) + BN + GELU\n"
    "  Conv2d(256→256, k3, s1, p1) + BN + GELU\n"
    "  → output: (B, 64, 256)\n"
    "  # 64 patches, each 256-dim",
    font_size=11,
)

add_textbox(slide, 17.5, 10.5, 15, 0.6, "设计优势", font_size=13, bold=True, color=ACCENT)
tf = add_textbox(slide, 17.5, 11.2, 15, 6, "", font_size=12)
add_bullet_list(
    tf,
    [
        "4 层卷积逐步下采样，stride=2 两次使空间尺寸从 32→16→8，形成 8×8=64 个 patch",
        "卷积操作内置局部归纳偏置，缓解纯 ViT 在小数据集上的过拟合",
        "相比标准 PatchEmbed (224×224, 16×16 patch → 196 tokens)，tokens 数从 196 降至 64",
        "注意力计算量 O(N²) 大幅降低，训练速度显著提升",
    ],
    font_size=12,
)

add_page_number(slide, 4, TOTAL_PAGES)

# ===========================================================================
# Page 5: 核心模块实现 - Multi-Head Self-Attention
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "核心模块实现：Multi-Head Self-Attention", "注意力机制的前向传播过程")

# Left - Explanation
add_rect(slide, 1.0, 2.2, 14.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14, 0.6, "计算流程", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 3.1, 14.5, 6, "", font_size=12)
add_bullet_list(
    tf,
    [
        "输入 x: (B, N, C)，N=65 (1 CLS + 64 patches)",
        "通过 Linear 投影生成 Q、K、V 三个矩阵",
        "qkv = Linear(C → 3C)，reshape 为 (3, B, num_heads, N, head_dim)",
        "使用 scaled_dot_product_attention 计算注意力",
        "Attention(Q,K,V) = softmax(QK^T/√d_k) × V",
        "合并多头输出：transpose + reshape → (B, N, C)",
        "最终通过 Linear 投影回 dim 维度",
        "采用 PyTorch 2.x 的 F.scaled_dot_product_attention\n  自动使用 Flash Attention 加速",
    ],
    font_size=12,
)

# Right - Code
add_rect(slide, 16.5, 2.2, 16.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 16.7, 2.2, 15.5, 0.6, "核心代码", font_size=16, bold=True, color=WHITE)

add_code_block(
    slide,
    16.5,
    3.1,
    16.0,
    14.5,
    "class Attention(nn.Module):\n"
    "    def __init__(self, dim, num_heads=8,\n"
    "                 qkv_bias=False):\n"
    "        super().__init__()\n"
    "        self.num_heads = num_heads\n"
    "        head_dim = dim // num_heads\n"
    "        self.scale = head_dim ** -0.5\n"
    "        # QKV联合投影\n"
    "        self.qkv = nn.Linear(\n"
    "            dim, dim*3, bias=qkv_bias)\n"
    "        self.proj = nn.Linear(dim, dim)\n"
    "\n"
    "    def forward(self, x):\n"
    "        B, N, C = x.shape\n"
    "        qkv = self.qkv(x)\n"
    "        qkv = qkv.reshape(\n"
    "            B, N, 3, self.num_heads,\n"
    "            C//self.num_heads)\n"
    "        qkv = qkv.permute(\n"
    "            2, 0, 3, 1, 4)\n"
    "        q, k, v = qkv.unbind(0)\n"
    "        # Flash Attention\n"
    "        x = F.scaled_dot_product_attention(\n"
    "            q, k, v,\n"
    "            scale=self.scale)\n"
    "        x = x.transpose(1,2)\n"
    "        x = x.reshape(B, N, C)\n"
    "        return self.proj(x)",
    font_size=9,
)

add_page_number(slide, 5, TOTAL_PAGES)

# ===========================================================================
# Page 6: 核心模块实现 - Transformer Encoder & MLP
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "核心模块实现：Transformer Encoder & MLP", "编码器块结构与前馈网络")

# Left column - Block
add_rect(slide, 1.0, 2.2, 15.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 15, 0.6, "Transformer Encoder Block", font_size=14, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 3.0, 15, 2.5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "Pre-Norm 结构：LayerNorm 置于 Attention/MLP 之前",
        "LayerNorm → Multi-Head Self-Attention → DropPath → Residual (+)",
        "LayerNorm → MLP (GELU) → DropPath → Residual (+)",
        "Stochastic Depth：随机丢弃整个残差分支，增强泛化",
    ],
    font_size=12,
)

# Block code
add_code_block(
    slide,
    1.2,
    5.5,
    15.5,
    8.5,
    "class Block(nn.Module):\n"
    '    """Pre-Norm Transformer Block"""\n'
    "    def __init__(self, dim, num_heads,\n"
    "                 mlp_ratio=4., drop=0.,\n"
    "                 drop_path=0.):\n"
    "        super().__init__()\n"
    "        self.norm1 = LayerNorm(dim)\n"
    "        self.attn = Attention(\n"
    "            dim, num_heads)\n"
    "        self.norm2 = LayerNorm(dim)\n"
    "        self.mlp = Mlp(dim,\n"
    "            hidden=int(dim*mlp_ratio))\n"
    "        self.drop_path = DropPath(\n"
    "            drop_path)\n"
    "\n"
    "    def forward(self, x):\n"
    "        # Attention + residual\n"
    "        x = x + self.drop_path(\n"
    "            self.attn(\n"
    "            self.norm1(x)))\n"
    "        # MLP + residual\n"
    "        x = x + self.drop_path(\n"
    "            self.mlp(\n"
    "            self.norm2(x)))\n"
    "        return x",
    font_size=9,
)

# Right column - MLP
add_rect(slide, 17.5, 2.2, 15.5, 0.6, fill_color=ACCENT)
add_textbox(slide, 17.7, 2.2, 15, 0.6, "MLP 前馈网络", font_size=14, bold=True, color=WHITE)

tf = add_textbox(slide, 17.7, 3.0, 15, 2.5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "两层全连接 + GELU 激活 + Dropout",
        "隐藏层维度 = dim × mlp_ratio (256×3=768)",
        "结构：FC(256→768) → GELU → Drop → FC(768→256) → Drop",
        "GELU 相比 ReLU 更平滑，适合 Transformer",
    ],
    font_size=12,
)

add_code_block(
    slide,
    17.5,
    5.5,
    15.5,
    5.5,
    "class Mlp(nn.Module):\n"
    "    def __init__(self, in_features,\n"
    "                 hidden_features=None,\n"
    "                 drop=0.):\n"
    "        super().__init__()\n"
    "        hidden = hidden_features or in_features\n"
    "        self.fc1 = nn.Linear(in_features, hidden)\n"
    "        self.act = nn.GELU()\n"
    "        self.fc2 = nn.Linear(hidden, in_features)\n"
    "        self.drop = nn.Dropout(drop)\n"
    "\n"
    "    def forward(self, x):\n"
    "        x = self.fc1(x)\n"
    "        x = self.act(x)\n"
    "        x = self.drop(x)\n"
    "        x = self.fc2(x)\n"
    "        x = self.drop(x)\n"
    "        return x",
    font_size=9,
)

# Bottom - ViT forward
add_textbox(slide, 17.7, 11.5, 15, 0.6, "ViT 前向传播", font_size=13, bold=True, color=ACCENT)
add_code_block(
    slide,
    17.5,
    12.2,
    15.5,
    5.5,
    "class VisionTransformer(nn.Module):\n"
    "    def forward(self, x):\n"
    "        x = self.patch_embed(x)\n"
    "        # prepend CLS token + pos embed\n"
    "        cls = self.cls_token.expand(B,-1,-1)\n"
    "        x = torch.cat((cls, x), dim=1)\n"
    "        x = x + self.pos_embed\n"
    "        x = self.pos_drop(x)\n"
    "        # transformer encoder\n"
    "        x = self.blocks(x)\n"
    "        x = self.norm(x)\n"
    "        # take CLS token for classification\n"
    "        x = x[:, 0]\n"
    "        return self.head(x)",
    font_size=9,
)

add_page_number(slide, 6, TOTAL_PAGES)

# ===========================================================================
# Page 7: 损失函数与优化器设计
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "损失函数与优化器设计", "Loss Function | Optimizer | Learning Rate Schedule")

# Loss
add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "损失函数", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 3.0, 15, 3.5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "CrossEntropyWithSoftTargets：同时支持硬标签和 MixUp 软标签",
        "硬标签：标准交叉熵 + Label Smoothing (ε=0.05)",
        "软标签 (MixUp 生成)：直接使用 soft targets 计算 KL 散度",
        "标签平滑公式：y' = (1-ε)·y_onehot + ε/K",
    ],
    font_size=12,
)

add_code_block(
    slide,
    1.2,
    7.0,
    15,
    4,
    "class CrossEntropyWithSoftTargets(nn.Module):\n"
    "    def forward(self, logits, targets):\n"
    "        if targets.ndim == 1:\n"
    "            # 硬标签 + label smoothing\n"
    "            return F.cross_entropy(\n"
    "                logits, targets,\n"
    "                label_smoothing=0.05)\n"
    "        # MixUp 软标签\n"
    "        return -(targets\n"
    "                 * F.log_softmax(logits,dim=-1)\n"
    "                ).sum(dim=-1).mean()",
    font_size=9,
)

# Optimizer
add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "优化器设计", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 17.7, 3.0, 15, 3.5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "AdamW 优化器，lr=8e-4，weight_decay=0.05",
        "参数分组策略：",
        "  - Decay 参数：权重矩阵 (weight)",
        "  - No Decay 参数：bias、LayerNorm、CLS token、Position Embedding",
        "解耦权重衰减，避免对归一化层参数施加正则化",
    ],
    font_size=12,
)

add_code_block(
    slide,
    17.5,
    7.0,
    15,
    4,
    "# Separate param groups\n"
    "decay, no_decay = [], []\n"
    "for name, p in model.named_parameters():\n"
    "    if any(k in name for k in\n"
    '           ["bias","norm","cls_token",\n'
    '            "pos_embed"]):\n'
    "        no_decay.append(p)\n"
    "    else:\n"
    "        decay.append(p)\n"
    "optimizer = AdamW([\n"
    '    {"params": decay},\n'
    '    {"params": no_decay, "weight_decay": 0}\n'
    "], lr=8e-4, weight_decay=0.05)",
    font_size=9,
)

# LR Schedule
add_rect(slide, 1.0, 11.8, 31.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 11.8, 31, 0.6, "学习率调度策略", font_size=14, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 12.6, 31, 5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "Warmup (5 epochs)：LinearLR，start_factor=0.1，从小学习率线性增长，避免训练初期不稳定",
        "Cosine Annealing (剩余 55 epochs)：余弦退火衰减，平滑地将学习率降至接近 0",
        "实现方式：SequentialLR 组合 LinearLR + CosineAnnealingLR，自动衔接",
    ],
    font_size=12,
)

add_page_number(slide, 7, TOTAL_PAGES)

# ===========================================================================
# Page 8: 训练配置与创新点
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "训练配置与创新点", "Training Configuration & Design Innovations")

# Training config table
add_textbox(slide, 1.2, 2.2, 31, 0.6, "训练超参数", font_size=14, bold=True, color=PRIMARY)
add_table(
    slide,
    1.2,
    2.9,
    ["参数", "值", "参数", "值"],
    [
        ["Batch Size", "512", "Max Epochs", "60 (early stop @20)"],
        ["Learning Rate", "8e-4", "Weight Decay", "0.05"],
        ["Warmup Epochs", "5", "Label Smoothing", "0.05"],
        ["MixUp Alpha", "0.2", "Random Erasing p", "0.15"],
        ["Mixed Precision", "AMP (bfloat16)", "Optimizer", "AdamW"],
        ["Scheduler", "Warmup + Cosine", "Val Fraction", "0.1"],
    ],
    col_widths=[7, 8, 7, 8],
)

# Innovation points
add_textbox(slide, 1.2, 8.8, 31, 0.6, "创新设计点", font_size=14, bold=True, color=PRIMARY)

# Innovation 1
add_rect(slide, 1.2, 9.5, 15, 0.5, fill_color=ACCENT)
add_textbox(
    slide, 1.4, 9.5, 14.5, 0.5, "创新 1：卷积 Tokenizer 替代标准 Patch Embed", font_size=13, bold=True, color=WHITE
)
tf = add_textbox(slide, 1.4, 10.2, 15, 3, "", font_size=11)
add_bullet_list(
    tf,
    [
        "标准 ViT 将 32×32 放大到 224×224 → 196 tokens",
        "本方案用 4 层卷积直接处理 32×32 → 64 tokens",
        "tokens 减少 67%，注意力计算量降低 56%",
        "卷积提供局部归纳偏置，缓解小数据过拟合",
    ],
    font_size=11,
    bullet_char="·",
)

# Innovation 2
add_rect(slide, 17.5, 9.5, 15, 0.5, fill_color=ACCENT)
add_textbox(slide, 17.7, 9.5, 14.5, 0.5, "创新 2：综合正则化策略", font_size=13, bold=True, color=WHITE)
tf = add_textbox(slide, 17.7, 10.2, 15, 3, "", font_size=11)
add_bullet_list(
    tf,
    [
        "MixUp + Label Smoothing：软标签平滑，防止过拟合",
        "RandAugment + RandomErasing：多样化数据增强",
        "Stochastic Depth：随机深度正则化",
        "Weight Decay 分离：不衰减 bias/norm 参数",
    ],
    font_size=11,
    bullet_char="·",
)

# Innovation 3
add_rect(slide, 1.2, 13.5, 15, 0.5, fill_color=ACCENT)
add_textbox(slide, 1.4, 13.5, 14.5, 0.5, "创新 3：现代训练技术栈", font_size=13, bold=True, color=WHITE)
tf = add_textbox(slide, 1.4, 14.2, 15, 3, "", font_size=11)
add_bullet_list(
    tf,
    [
        "混合精度 AMP (bfloat16)：减少显存，加速训练",
        "torch.compile (reduce-overhead)：JIT 编译优化",
        "channels_last 内存格式：提升 CUDA 吞吐",
        "F.scaled_dot_product_attention：自动启用 Flash Attention",
    ],
    font_size=11,
    bullet_char="·",
)

# Innovation 4
add_rect(slide, 17.5, 13.5, 15, 0.5, fill_color=ACCENT)
add_textbox(slide, 17.7, 13.5, 14.5, 0.5, "创新 4：提前停止与验证策略", font_size=13, bold=True, color=WHITE)
tf = add_textbox(slide, 17.7, 14.2, 15, 3, "", font_size=11)
add_bullet_list(
    tf,
    [
        "从训练集划分 10% 作为独立验证集",
        "基于验证准确率保存最佳模型 (best_model.pt)",
        "达到目标准确率 (82%) 后提前停止",
        "最少训练 min_epochs=20 保证充分学习",
    ],
    font_size=11,
    bullet_char="·",
)

add_page_number(slide, 8, TOTAL_PAGES)

# ===========================================================================
# Page 9: 实验结果与分析
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "实验结果与分析", "训练过程 | 最终指标 | 训练曲线 | 混淆矩阵")

# Final results highlight
add_rect(slide, 1.0, 2.2, 31.5, 1.5, fill_color=LIGHT_BG)
add_textbox(
    slide,
    2.0,
    2.3,
    6.5,
    1.3,
    f"最佳验证准确率\n{best_val_acc:.2%}",
    font_size=24,
    bold=True,
    color=PRIMARY,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    10.0,
    2.3,
    6.5,
    1.3,
    f"测试准确率\n{test_acc:.2%}",
    font_size=24,
    bold=True,
    color=ACCENT,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    18.5,
    2.3,
    6.5,
    1.3,
    f"测试损失\n{test_loss:.4f}",
    font_size=24,
    bold=True,
    color=RGBColor(0xE6, 0x7E, 0x22),
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    26,
    2.3,
    6.5,
    1.3,
    "达标 Epoch\n12",
    font_size=24,
    bold=True,
    color=RGBColor(0x27, 0xAE, 0x60),
    alignment=PP_ALIGN.CENTER,
)

# Training curves image
add_textbox(slide, 1.5, 4.2, 15, 0.6, "训练曲线", font_size=14, bold=True, color=PRIMARY)
add_image(slide, OUTPUT_DIR / "training_curves.png", 1.5, 4.8, 15, 6.5)

# Confusion matrix image
add_textbox(slide, 17.5, 4.2, 15, 0.6, "混淆矩阵", font_size=14, bold=True, color=PRIMARY)
add_image(slide, OUTPUT_DIR / "confusion_matrix.png", 17.5, 4.8, 15, 6.5)

# Analysis
add_rect(slide, 1.0, 11.8, 31.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 11.8, 31, 0.6, "结果分析", font_size=14, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 12.6, 31, 5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "第 12 epoch 验证准确率已达 80.02%，第 20 epoch 达到最佳 84.60%，收敛速度快",
        "训练准确率低于验证准确率 (MixUp 效应)，评估模型应参考验证集和测试集指标",
        "32×32 原生输入 + 卷积 tokenizer 有效减少 tokens 数量和注意力计算量",
        "混淆矩阵显示大多数类别分类准确，少数混淆集中在相似类别间 (如 cat↔dog, deer↔horse)",
    ],
    font_size=12,
)

add_page_number(slide, 9, TOTAL_PAGES)

# ===========================================================================
# Page 10: 预测样本、特征可视化与总结
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "预测样本、特征可视化与总结", "Predictions | Feature Maps | Summary")

# Predictions image
add_textbox(slide, 1.5, 2.2, 15, 0.6, "测试集预测样本", font_size=14, bold=True, color=PRIMARY)
add_image(slide, OUTPUT_DIR / "predictions.png", 1.5, 2.8, 15, 6.0)

# Feature maps image
add_textbox(slide, 17.5, 2.2, 15, 0.6, "第一层卷积特征图", font_size=14, bold=True, color=PRIMARY)
add_image(slide, OUTPUT_DIR / "feature_maps.png", 17.5, 2.8, 15, 6.0)

# Summary
add_rect(slide, 1.0, 9.2, 31.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 9.2, 31, 0.6, "总结", font_size=16, bold=True, color=WHITE)

add_textbox(slide, 1.5, 10.0, 15, 0.6, "已完成工作", font_size=13, bold=True, color=ACCENT)
tf = add_textbox(slide, 1.5, 10.6, 15, 3, "", font_size=12)
add_bullet_list(
    tf,
    [
        "从零实现 ViT 所有核心模块 (Attention, Block, MLP, PatchEmbed)",
        "完成 CIFAR-10 数据加载、增强、训练、验证、测试全流程",
        "测试准确率 83.36%，超过实验要求 (80%)",
        "生成训练曲线、混淆矩阵、预测样本和特征图可视化",
    ],
    font_size=12,
)

add_textbox(slide, 17.5, 10.0, 15, 0.6, "改进方向", font_size=13, bold=True, color=ACCENT)
tf = add_textbox(slide, 17.5, 10.6, 15, 3, "", font_size=12)
add_bullet_list(
    tf,
    [
        "延长训练至更多 epoch，配合更积极的余弦衰减",
        "采用 CutMix + MixUp 联合增强策略",
        "使用预训练 ViT 权重进行迁移学习或知识蒸馏",
        "尝试更大的模型容量 (更多层数/更宽嵌入维度)",
    ],
    font_size=12,
)

add_page_number(slide, 10, TOTAL_PAGES)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(PPT_PATH))
print(f"PPT saved to: {PPT_PATH}")
print(f"Total slides: {len(prs.slides)}")
