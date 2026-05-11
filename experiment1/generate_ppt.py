"""
Generate experiment1 presentation PPT based on the report outline.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
OUTPUT_DIR = Path(__file__).resolve().parent / "outputs"
PPT_PATH = OUTPUT_DIR / "experiment1_report.pptx"
METRICS_PATH = OUTPUT_DIR / "metrics.json"

# ---------------------------------------------------------------------------
# Theme colors
# ---------------------------------------------------------------------------
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)
SECONDARY = RGBColor(0x2D, 0x3A, 0x4A)
ACCENT = RGBColor(0x00, 0x96, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFA)
DARK_TEXT = RGBColor(0x1E, 0x1E, 0x1E)
CODE_BG = RGBColor(0x28, 0x2C, 0x34)
CODE_FG = RGBColor(0xAB, 0xB2, 0xBF)
TABLE_HEADER_BG = RGBColor(0x1A, 0x56, 0xDB)
TABLE_ROW_ALT = RGBColor(0xE8, 0xEE, 0xF8)

SLIDE_W = Cm(33.867)  # 16:9
SLIDE_H = Cm(19.05)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    bold=False,
    color=DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    font_name="Microsoft YaHei",
):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Cm(left),
        Cm(top),
        Cm(width),
        Cm(height),
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_bullet_list(tf, items, font_size=14, color=DARK_TEXT, bullet_char="●"):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(6)


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = add_rect(slide, left, top, width, height, fill_color=CODE_BG)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.4)
    tf.margin_right = Cm(0.4)
    tf.margin_top = Cm(0.3)
    tf.margin_bottom = Cm(0.3)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = CODE_FG
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.LEFT


def add_image(slide, path, left, top, width, height=None):
    if height is None:
        slide.shapes.add_picture(str(path), Cm(left), Cm(top), Cm(width))
    else:
        slide.shapes.add_picture(str(path), Cm(left), Cm(top), Cm(width), Cm(height))


def add_page_number(slide, num, total):
    add_textbox(
        slide,
        30,
        17.8,
        3.5,
        0.8,
        f"{num} / {total}",
        font_size=10,
        color=RGBColor(0x99, 0xAA, 0xBB),
        alignment=PP_ALIGN.RIGHT,
    )


def add_section_header(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 33.867, 1.0, fill_color=PRIMARY)
    add_textbox(slide, 1.0, 0.1, 30, 0.9, title, font_size=22, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 1.0, 1.15, 30, 0.8, subtitle, font_size=12, color=RGBColor(0x88, 0x99, 0xAA))


def add_table(slide, left, top, headers, rows, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Cm(left), Cm(top), Cm(sum(col_widths or [6] * n_cols)), Cm(0.8 * n_rows)
    )
    tbl = table_shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER

    return table_shape


# ---------------------------------------------------------------------------
# Load metrics
# ---------------------------------------------------------------------------
metrics = json.loads(METRICS_PATH.read_text())
best_val_acc = metrics["best_validation_accuracy"]
test_acc = metrics["test_accuracy"]
test_loss = metrics["test_loss"]

TOTAL_PAGES = 11

# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

# ===========================================================================
# Page 1: 封面
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)

add_rect(slide, 0, 0, 33.867, 0.3, fill_color=PRIMARY)

add_rect(slide, 0, 3.5, 33.867, 9.0, fill_color=LIGHT_BG)
add_textbox(
    slide, 2.0, 4.5, 30, 2.0, "实验一：手写数字识别", font_size=36, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER
)
add_textbox(
    slide, 2.0, 7.0, 30, 1.5, "基于 CNN 的 MNIST 图像分类", font_size=20, color=SECONDARY, alignment=PP_ALIGN.CENTER
)

add_textbox(
    slide,
    2.0,
    13.5,
    30,
    0.8,
    "深度学习 · 2026 Spring",
    font_size=14,
    color=RGBColor(0x88, 0x99, 0xAA),
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    2.0,
    14.5,
    30,
    0.8,
    "李相廷 · 2026-04-15",
    font_size=12,
    color=RGBColor(0x99, 0xAA, 0xBB),
    alignment=PP_ALIGN.CENTER,
)

add_page_number(slide, 1, TOTAL_PAGES)

# ===========================================================================
# Page 2: 目录
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "目录", "Contents")

toc_items = [
    "1. 实验任务概述",
    "2. 数据集介绍",
    "3. 模型结构设计",
    "4. 损失函数与优化器",
    "5. 创新点说明",
    "6. 训练过程与配置",
    "7. 实验结果",
    "8. 结果分析与讨论",
    "9. 总结与改进方向",
]

for i, item in enumerate(toc_items):
    y = 2.5 + i * 1.5
    add_rect(slide, 4.0, y, 0.4, 0.4, fill_color=PRIMARY)
    add_textbox(slide, 5.0, y - 0.1, 25, 1.0, item, font_size=18, color=SECONDARY)

add_page_number(slide, 2, TOTAL_PAGES)

# ===========================================================================
# Page 3: 实验任务概述
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "实验任务概述", "任务定义 | 应用场景 | 实验目标")

add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "任务定义", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 3.0, 15.0, 4.5, "", font_size=13)
add_bullet_list(
    tf,
    [
        "对手写数字灰度图像进行 0-9 十个类别的分类",
        "输入：28×28 单通道灰度手写数字图像",
        "输出：0-9 的类别标签",
        "属于典型的监督式图像分类任务",
    ],
    font_size=13,
)

add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "应用场景", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 17.7, 3.0, 15.0, 4.5, "", font_size=13)
add_bullet_list(
    tf,
    [
        "邮政编码自动识别",
        "票据金额数字化录入",
        "银行表单自动化处理",
        "各类文档中的手写数字 OCR",
    ],
    font_size=13,
)

add_rect(slide, 1.0, 8.0, 31.5, 0.6, fill_color=ACCENT)
add_textbox(slide, 1.2, 8.0, 31, 0.6, "实验目标", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 8.9, 31, 8.0, "", font_size=13)
add_bullet_list(
    tf,
    [
        "搭建基于 PyTorch 的 CNN 卷积神经网络模型",
        "在 MNIST 训练集（60,000 张）上完成模型训练",
        "在 MNIST 测试集（10,000 张）上达到 98% 以上分类准确率",
        "使用 Python-PPTX 生成实验汇报 PPT",
    ],
    font_size=13,
)

add_page_number(slide, 3, TOTAL_PAGES)

# ===========================================================================
# Page 4: 数据集介绍
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "数据集介绍", "MNIST 手写数字数据库详情与预处理")

add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "数据集详情", font_size=16, bold=True, color=WHITE)

add_table(
    slide,
    1.2,
    3.2,
    ["属性", "详情"],
    [
        ["来源", "MNIST (LeCun et al., 1998)"],
        ["训练集规模", "60,000 张 → 训练 55K + 验证 5K"],
        ["测试集规模", "10,000 张"],
        ["类别", "10 类（数字 0-9），分布均衡"],
        ["图像尺寸", "28×28 单通道灰度图"],
    ],
    col_widths=[7, 8],
)

add_textbox(slide, 1.2, 8.2, 15, 0.6, "10 个类别", font_size=14, bold=True, color=SECONDARY)
tf = add_textbox(slide, 1.2, 8.8, 15, 3, "", font_size=13)
add_bullet_list(tf, ["0 · 1 · 2 · 3 · 4 · 5 · 6 · 7 · 8 · 9"])

add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "数据预处理与增强", font_size=16, bold=True, color=WHITE)

add_textbox(slide, 17.7, 3.2, 14.5, 0.6, "标准化", font_size=13, bold=True, color=ACCENT)
tf = add_textbox(slide, 17.7, 3.8, 14.5, 1.5, "", font_size=11)
add_bullet_list(
    tf,
    [
        "Normalize(mean=0.1307, std=0.3081)",
        "基于 MNIST 全集的统计量",
    ],
    font_size=11,
    bullet_char="·",
)

add_textbox(slide, 17.7, 5.5, 14.5, 0.6, "训练数据增强（仅训练集）", font_size=13, bold=True, color=ACCENT)
tf = add_textbox(slide, 17.7, 6.1, 14.5, 4, "", font_size=11)
add_bullet_list(
    tf,
    [
        "RandomAffine(degrees=±12°, translate=(0.08, 0.08))",
        "增加模型对旋转和平移的鲁棒性",
        "验证集和测试集仅做 ToTensor + Normalize",
        "确保评估公平性",
    ],
    font_size=11,
)

add_textbox(slide, 17.7, 10.5, 14.5, 0.6, "数据划分策略", font_size=13, bold=True, color=ACCENT)
tf = add_textbox(slide, 17.7, 11.1, 14.5, 4, "", font_size=11)
add_bullet_list(
    tf,
    [
        "训练集：前 55,000 张（带数据增强）",
        "验证集：后 5,000 张（仅标准化）",
        "测试集：10,000 张（仅标准化）",
        "固定随机种子 (seed=42)，确保可重复划分",
    ],
    font_size=11,
)

add_page_number(slide, 4, TOTAL_PAGES)

# ===========================================================================
# Page 5: 模型结构设计
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "模型结构设计", "CNN 卷积神经网络架构")

add_rect(slide, 1.0, 2.2, 31.5, 2.2, fill_color=LIGHT_BG)
add_textbox(slide, 1.5, 2.3, 30.5, 0.6, "整体流程", font_size=14, bold=True, color=PRIMARY)
add_textbox(
    slide,
    2.0,
    3.0,
    30,
    1.2,
    "Input(1×28×28)  →  Conv-BN-ReLU ×2  →  MaxPool + Dropout  →  "
    "Conv-BN-ReLU ×2  →  MaxPool + Dropout  →  Flatten  →  FC(3136→128)  →  ReLU + Dropout  →  FC(128→10)",
    font_size=13,
    color=SECONDARY,
    alignment=PP_ALIGN.CENTER,
)

add_textbox(slide, 1.5, 4.8, 31, 0.6, "网络结构参数", font_size=14, bold=True, color=PRIMARY)
add_table(
    slide,
    1.5,
    5.5,
    ["模块", "配置", "输出尺寸"],
    [
        ["Input", "—", "1 × 28 × 28"],
        ["Conv2d(1→32, k3) + BN + ReLU", "×2", "32 × 28 × 28"],
        ["MaxPool2d(k2) + Dropout(0.25)", "—", "32 × 14 × 14"],
        ["Conv2d(32→64, k3) + BN + ReLU", "×2", "64 × 14 × 14"],
        ["MaxPool2d(k2) + Dropout(0.25)", "—", "64 × 7 × 7"],
        ["Flatten + Linear(3136→128) + ReLU + Dropout(0.5)", "—", "128"],
        ["Linear(128→10)", "—", "10 (logits)"],
    ],
    col_widths=[12, 6, 10],
)

add_textbox(slide, 1.5, 12.2, 31, 0.6, "设计特点", font_size=14, bold=True, color=PRIMARY)
tf = add_textbox(slide, 1.5, 12.9, 31, 5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "通道数逐步增加（1→32→64），从底层边缘特征到高层语义特征逐步抽象",
        "MaxPool 两次降采样（28→14→7），逐步压缩空间维度，扩大感受野",
        "BatchNorm 加速收敛并稳定训练，Dropout 在不同位置使用不同比例防止过拟合",
        "参数量约 44 万，属于轻量级模型，适合 MNIST 任务规模",
    ],
    font_size=12,
)

add_page_number(slide, 5, TOTAL_PAGES)

# ===========================================================================
# Page 6: 损失函数与优化器
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "损失函数与优化器", "Loss Function | Optimizer | AMP 混合精度")

# Loss
add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "损失函数", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 3.0, 15, 5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "CrossEntropyLoss（交叉熵损失）",
        "适用于多分类任务的标准损失函数",
        "PyTorch 内部自动融合 softmax + NLLLoss",
        "数值更稳定，避免单独计算 softmax 的数值问题",
        "公式：L = -Σ y_true · log(y_pred)",
    ],
    font_size=12,
)

add_code_block(slide, 1.2, 8.5, 15, 1.5, "criterion = nn.CrossEntropyLoss()", font_size=11)

# Optimizer
add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "优化器设计", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 17.7, 3.0, 15, 5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "Adam 优化器",
        "lr = 0.001（1e-3）",
        "weight_decay = 1e-4（L2 正则化）",
        "结合动量与自适应学习率",
        "收敛速度快，超参数不敏感",
        "适合中小规模实验",
    ],
    font_size=12,
)

add_code_block(
    slide,
    17.5,
    8.5,
    15,
    2,
    "optimizer = torch.optim.Adam(\n    model.parameters(),\n    lr=1e-3,\n    weight_decay=1e-4\n)",
    font_size=10,
)

# AMP
add_rect(slide, 1.0, 10.8, 31.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 10.8, 31, 0.6, "AMP 混合精度训练", font_size=14, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 11.6, 31, 5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "使用 torch.amp.autocast + GradScaler 实现自动混合精度",
        "CUDA 可用时自动启用，在 RTX 5070 上利用 BF16 加速训练",
        "训练效率提升约 1.5×，同时减少显存占用",
        "通过早停策略保存验证准确率最优的模型 checkpoint",
    ],
    font_size=12,
)

add_page_number(slide, 6, TOTAL_PAGES)

# ===========================================================================
# Page 7: 创新点说明
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "创新点说明", "Design Innovations")

# Innovation 1
add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "创新 1：数据增强策略", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 1.2, 3.0, 15, 6, "", font_size=12)
add_bullet_list(
    tf,
    [
        "RandomAffine：±12° 旋转 + ±8% 平移 + ±5% 缩放",
        "仅作用于训练集，验证集和测试集保持纯净",
        "增加模型对平移、旋转、缩放的鲁棒性",
        "模拟手写数字在实际场景中的变形",
        "不增加推理阶段计算开销",
    ],
    font_size=12,
)

# Innovation 2
add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "创新 2：Dropout 分层正则化", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 17.7, 3.0, 15, 6, "", font_size=12)
add_bullet_list(
    tf,
    [
        "特征提取器中使用 Dropout(p=0.25)：轻量正则化，保留更多底层特征",
        "分类器中使用 Dropout(p=0.5)：强正则化，防止全连接层过拟合",
        "训练时生效，推理时自动关闭",
        "训练 accuracy 始终低于验证 accuracy，证明 Dropout 有效",
        "训练/验证 loss 差距小，泛化能力强",
    ],
    font_size=12,
)

# Innovation 3
add_rect(slide, 1.0, 9.5, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 1.2, 9.5, 14.5, 0.6, "创新 3：AMP 混合精度训练", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 1.2, 10.3, 15, 5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "torch.amp.autocast 自动管理精度转换",
        "GradScaler 防止低精度下梯度下溢",
        "在 RTX 5070 上实现约 1.5× 加速",
        "最终精度无损失（测试准确率 99.51%）",
        "显存占用减少，可支持更大 batch size",
    ],
    font_size=12,
)

# Innovation 4
add_rect(slide, 17.5, 9.5, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 17.7, 9.5, 14.5, 0.6, "创新 4：早停与模型选择", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 17.7, 10.3, 15, 5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "每 epoch 验证后记录并比较验证准确率",
        "仅保存验证集最优 checkpoint（best_model.pt）",
        "保证最终使用的是泛化能力最强的模型",
        "避免过拟合阶段的模型被用于测试",
        "最终测试使用独立的测试集，确保评估无偏",
    ],
    font_size=12,
)

add_page_number(slide, 7, TOTAL_PAGES)

# ===========================================================================
# Page 8: 训练过程与配置
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "训练过程与配置", "Training Configuration & Pipeline")

add_textbox(slide, 1.2, 2.2, 31, 0.6, "训练超参数", font_size=14, bold=True, color=PRIMARY)
add_table(
    slide,
    1.2,
    2.9,
    ["参数", "值", "参数", "值"],
    [
        ["Batch Size", "128", "Max Epochs", "8"],
        ["Learning Rate", "0.001", "Weight Decay", "1e-4"],
        ["Optimizer", "Adam", "Loss", "CrossEntropyLoss"],
        ["AMP", "Autocast + GradScaler", "Seed", "42"],
        ["Augmentation", "RandomAffine (±12°, ±8%)", "Val Size", "5,000"],
        ["Hardware", "NVIDIA RTX 5070", "Framework", "PyTorch 2.11"],
    ],
    col_widths=[7, 8, 7, 8],
)

add_textbox(slide, 1.2, 8.6, 31, 0.6, "训练流程", font_size=14, bold=True, color=PRIMARY)
tf = add_textbox(slide, 1.2, 9.3, 31, 8, "", font_size=12)
add_bullet_list(
    tf,
    [
        "Step 1：加载 MNIST 数据集，划分训练集 55K + 验证集 5K，应用数据增强和标准化",
        "Step 2：构建 MNISTCNN 模型，约 44 万参数，部署到 CUDA 设备",
        "Step 3：8 轮训练/验证循环，每轮记录 train_loss、train_acc、val_loss、val_acc",
        "Step 4：自动保存验证准确率最高的模型 checkpoint（best_model.pt）",
        "Step 5：加载最优模型，在独立测试集（10,000 张）上评估",
        "Step 6：输出 metrics.json（含完整训练历史和最终指标）",
    ],
    font_size=12,
)

# Epoch results
add_textbox(slide, 1.2, 14.2, 31, 0.6, "各 Epoch 验证准确率", font_size=13, bold=True, color=ACCENT)
add_table(
    slide,
    1.2,
    14.9,
    ["Epoch", "1", "2", "3", "4", "5", "6", "7", "8"],
    [["Val Acc", "98.04%", "98.46%", "98.78%", "98.94%", "99.04%", "99.04%", "99.20%", "99.36%"]],
    col_widths=[3] + [2.5] * 8,
)

add_page_number(slide, 8, TOTAL_PAGES)

# ===========================================================================
# Page 9: 实验结果
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "实验结果", "关键指标 | 训练曲线 | 混淆矩阵")

# Key metrics highlight
add_rect(slide, 1.0, 2.2, 31.5, 1.5, fill_color=LIGHT_BG)
add_textbox(
    slide,
    2.0,
    2.3,
    7.5,
    1.3,
    f"最佳验证准确率\n{best_val_acc:.2%}",
    font_size=24,
    bold=True,
    color=PRIMARY,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    11.0,
    2.3,
    7.5,
    1.3,
    f"测试准确率\n{test_acc:.2%}",
    font_size=24,
    bold=True,
    color=ACCENT,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    20.0,
    2.3,
    7.5,
    1.3,
    f"测试损失\n{test_loss:.4f}",
    font_size=24,
    bold=True,
    color=RGBColor(0xE6, 0x7E, 0x22),
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    29.0,
    2.3,
    4.5,
    1.3,
    "Epochs\n8",
    font_size=24,
    bold=True,
    color=RGBColor(0x27, 0xAE, 0x60),
    alignment=PP_ALIGN.CENTER,
)

# Training curves
add_textbox(slide, 1.5, 4.2, 15, 0.6, "训练曲线", font_size=14, bold=True, color=PRIMARY)
add_image(slide, OUTPUT_DIR / "training_curves.png", 1.5, 4.8, 15, 6.0)

# Confusion matrix
add_textbox(slide, 17.5, 4.2, 15, 0.6, "混淆矩阵", font_size=14, bold=True, color=PRIMARY)
add_image(slide, OUTPUT_DIR / "confusion_matrix.png", 17.5, 4.8, 15, 6.0)

# Observations
add_rect(slide, 1.0, 11.3, 31.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 11.3, 31, 0.6, "关键发现", font_size=14, bold=True, color=WHITE)

tf = add_textbox(slide, 1.2, 12.1, 31, 5, "", font_size=12)
add_bullet_list(
    tf,
    [
        "训练 loss 持续下降（0.488 → 0.109），验证准确率稳步提升（98.04% → 99.36%），未见明显过拟合",
        "混淆矩阵对角线主导，大部分类别准确率接近 100%；误差集中在形态相似类别（4↔9, 7↔2）",
        f"测试准确率 {test_acc:.2%}，远超实验要求的 98%，验证 CNN 对 MNIST 任务的高效性",
    ],
    font_size=12,
)

add_page_number(slide, 9, TOTAL_PAGES)

# ===========================================================================
# Page 10: 结果分析与讨论
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "结果分析与讨论", "Analysis & Discussion")

# Analysis 1
add_rect(slide, 1.0, 2.2, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 14.5, 0.6, "准确率超预期", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 1.2, 3.0, 15, 4, "", font_size=12)
add_bullet_list(
    tf,
    [
        f"测试准确率 {test_acc:.2%}，显著超过 98% 要求",
        "CNN 对 MNIST 这类规整灰度图像分类任务十分有效",
        "仅 8 个 epoch 即收敛至最优",
        "模型已充分学习手写数字的判别性特征",
    ],
    font_size=12,
)

# Analysis 2
add_rect(slide, 17.5, 2.2, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 17.7, 2.2, 14.5, 0.6, "泛化能力", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 17.7, 3.0, 15, 4, "", font_size=12)
add_bullet_list(
    tf,
    [
        "Dropout + 数据增强有效防止过拟合",
        "训练/验证 loss 差距小，模型泛化良好",
        "训练 accuracy 始终低于验证 accuracy",
        "Dropout 在训练阶段生效，验证阶段关闭，正确发挥了正则化作用",
    ],
    font_size=12,
)

# Analysis 3
add_rect(slide, 1.0, 7.0, 15.0, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 7.0, 14.5, 0.6, "特征学习", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 1.2, 7.8, 15, 4, "", font_size=12)
add_bullet_list(
    tf,
    [
        "Conv1 输出 32 个通道，可视化显示网络学习了有意义的底层视觉特征",
        "部分通道捕获笔画边缘方向（水平/垂直/斜向）",
        "部分通道捕获笔画粗细和轮廓信息",
        "特征层次合理：边缘检测器 → 形状模式 → 数字语义",
    ],
    font_size=12,
)

# Analysis 4
add_rect(slide, 17.5, 7.0, 15.0, 0.6, fill_color=ACCENT)
add_textbox(slide, 17.7, 7.0, 14.5, 0.6, "误差分析", font_size=14, bold=True, color=WHITE)
tf = add_textbox(slide, 17.7, 7.8, 15, 4, "", font_size=12)
add_bullet_list(
    tf,
    [
        "主要错误出现在形态相似的数字之间：4↔9, 7↔2",
        "符合人类认知规律——这些数字对确实容易混淆",
        "错误集中在笔迹潦草或形态模糊的样本上",
        "可通过更多数据增强或针对性样本训练改善",
    ],
    font_size=12,
)

# Prediction samples image
add_textbox(slide, 1.5, 11.8, 15, 0.6, "测试集预测样本（绿=正确，红=错误）", font_size=13, bold=True, color=PRIMARY)
add_image(slide, OUTPUT_DIR / "predictions.png", 1.5, 12.4, 15, 5.5)

# Feature maps image
add_textbox(slide, 17.5, 11.8, 15, 0.6, "第一层卷积特征图可视化", font_size=13, bold=True, color=PRIMARY)
add_image(slide, OUTPUT_DIR / "feature_maps.png", 17.5, 12.4, 15, 5.5)

add_page_number(slide, 10, TOTAL_PAGES)

# ===========================================================================
# Page 11: 总结与改进方向
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, WHITE)
add_section_header(slide, "总结与改进方向", "Summary & Future Work")

# Summary
add_rect(slide, 1.0, 2.2, 31.5, 0.6, fill_color=PRIMARY)
add_textbox(slide, 1.2, 2.2, 31, 0.6, "完成情况", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.5, 3.0, 31, 4, "", font_size=13)
add_bullet_list(
    tf,
    [
        f"实现了基于 PyTorch 的 CNN 模型，完成 MNIST 手写数字识别任务，测试准确率 {test_acc:.2%}",
        "完成从数据加载、模型构建、训练验证、到测试评估的完整深度学习流程",
        "训练曲线、混淆矩阵、预测样本、特征图等可视化全面展示了模型性能",
        "达成并超过实验指导书要求的 98% 准确率目标",
    ],
    font_size=13,
)

# Improvements
add_rect(slide, 1.0, 7.5, 31.5, 0.6, fill_color=ACCENT)
add_textbox(slide, 1.2, 7.5, 31, 0.6, "可改进方向", font_size=16, bold=True, color=WHITE)

tf = add_textbox(slide, 1.5, 8.3, 15, 8, "", font_size=12)
add_bullet_list(
    tf,
    [
        "引入学习率调度策略",
        "  - CosineAnnealingLR：余弦退火平滑衰减",
        "  - ReduceLROnPlateau：验证停滞时自动降低学习率",
        "尝试更深层网络结构",
        "  - 增加卷积层数或通道数",
        "  - 引入残差连接（ResNet 风格）",
    ],
    font_size=12,
)

tf = add_textbox(slide, 17.5, 8.3, 15, 8, "", font_size=12)
add_bullet_list(
    tf,
    [
        "增强策略升级",
        "  - RandomErasing：随机遮挡增强鲁棒性",
        "  - RandAugment：自动化增强组合搜索",
        "  - CutMix / MixUp：混合样本增强",
        "模型集成与蒸馏",
        "  - 多模型投票提升准确率",
        "  - 知识蒸馏压缩模型规模",
    ],
    font_size=12,
)

add_page_number(slide, 11, TOTAL_PAGES)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(PPT_PATH))
print(f"PPT saved to: {PPT_PATH}")
print(f"Total slides: {len(prs.slides)}")
